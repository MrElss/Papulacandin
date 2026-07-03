#!/usr/bin/env python3
"""
make_promotion_deck.py
======================
Build a professional, figure-rich promotion-report deck (16:9) for the
Papulacandin serum-gap program (Phases 1-14), for a general-faculty audience.

Embeds the cohesive figure set from make_deck_figures.py (run that first) plus a
real result plot; framework accents (chips, cards, roadmap loop, contributions
grid) are drawn with shapes.

Output: analysis/outputs/Papulacandin_promotion_report.pptx
"""
import os
import struct

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(HERE, "outputs")
FIG = os.path.join(OUT, "deck_figures")

INK = RGBColor(0x14, 0x23, 0x3A); BLUE = RGBColor(0x2F, 0x56, 0xA6)
TEAL = RGBColor(0x1F, 0x8A, 0x8A); CORAL = RGBColor(0xD8, 0x4A, 0x3C)
GREEN = RGBColor(0x2E, 0x8B, 0x57); AMBER = RGBColor(0xD8, 0x8A, 0x2B)
GREY = RGBColor(0x5B, 0x64, 0x72); LIGHT = RGBColor(0xF2, 0xF5, 0xF9)
CARD2 = RGBColor(0xE8, 0xEE, 0xF6); LINE = RGBColor(0xD3, 0xDA, 0xE4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
ML = Inches(0.75); CW = SW - 2 * ML
_page = 0


def _solid(sh, c): sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()
def _bg(s, c=WHITE): s.background.fill.solid(); s.background.fill.fore_color.rgb = c


def rect(s, l, t, w, h, c, shape=MSO_SHAPE.RECTANGLE, line=None, lw=1.0):
    sh = s.shapes.add_shape(shape, int(l), int(t), int(w), int(h)); _solid(sh, c)
    if line is not None:
        sh.line.color.rgb = line; sh.line.width = Pt(lw)
        sh.line.fill.solid(); sh.line.fill.fore_color.rgb = line
    sh.shadow.inherit = False
    return sh


def txt(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(int(l), int(t), int(w), int(h)); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2); tf.margin_top = tf.margin_bottom = Pt(1)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(4)
        for (text, size, bold, color, *rest) in para:
            r = p.add_run(); r.text = text; r.font.size = Pt(size)
            r.font.bold = bold; r.font.color.rgb = color; r.font.name = "Calibri"
            if rest and rest[0]:
                r.font.italic = True
    return tb


def chevron(s, l, t, w=Inches(0.34), h=Inches(0.34), c=GREY):
    rect(s, l, t, w, h, c, shape=MSO_SHAPE.CHEVRON)


def card(s, l, t, w, h, fill=LIGHT, line=LINE):
    return rect(s, l, t, w, h, fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=line)


def _png_size(path):
    with open(path, "rb") as f:
        head = f.read(26)
    return struct.unpack(">II", head[16:24])  # PNG width,height


def fig(s, name, l, t, boxw, boxh):
    """Fit-and-center an image within (l,t,boxw,boxh)."""
    path = os.path.join(FIG, name) if not os.path.isabs(name) else name
    if not os.path.exists(path):
        path2 = os.path.join(OUT, name)
        path = path2 if os.path.exists(path2) else path
    iw, ih = _png_size(path); ar = iw / ih; boxar = boxw / boxh
    if ar > boxar:
        w = boxw; h = boxw / ar
    else:
        h = boxh; w = boxh * ar
    x = l + (boxw - w) / 2; y = t + (boxh - h) / 2
    s.shapes.add_picture(path, int(x), int(y), width=int(w), height=int(h))


def base(s):
    _bg(s); rect(s, 0, 0, SW, Inches(0.16), INK)
    rect(s, 0, SH - Inches(0.16), SW, Inches(0.16), BLUE)


def header(kicker, title):
    global _page
    s = prs.slides.add_slide(BLANK); base(s)
    txt(s, ML, Inches(0.40), CW, Inches(0.30), [[(kicker.upper(), 13, True, BLUE)]])
    txt(s, ML, Inches(0.68), CW, Inches(0.70), [[(title, 26, True, INK)]])
    rect(s, ML, Inches(1.44), Inches(1.6), Pt(3), TEAL)
    _page += 1
    txt(s, SW - Inches(1.1), SH - Inches(0.5), Inches(0.8), Inches(0.3),
        [[(str(_page), 10, False, GREY)]], align=PP_ALIGN.RIGHT)
    return s


def strip(s, text_runs, top=Inches(6.35), fill=INK):
    card(s, ML, top, CW, Inches(0.9), fill, line=None)
    txt(s, ML + Inches(0.4), top + Inches(0.06), CW - Inches(0.8), Inches(0.78),
        text_runs, anchor=MSO_ANCHOR.MIDDLE)


# =========================================================== 1 TITLE
s = prs.slides.add_slide(BLANK); _bg(s, INK)
rect(s, 0, 0, SW, Inches(0.22), BLUE); rect(s, 0, SH - Inches(0.22), SW, Inches(0.22), TEAL)
txt(s, Inches(0.9), Inches(1.6), Inches(11.5), Inches(0.4),
    [[("PROMOTION REVIEW  ·  RESEARCH SUMMARY", 14, True, RGBColor(0x9C, 0xB4, 0xD8))]])
txt(s, Inches(0.9), Inches(2.15), Inches(11.6), Inches(2.4),
    [[("Designing Serum-Tolerant Antifungal Candidates", 40, True, WHITE)],
     [("An AI- and Quantum-Chemistry–Guided Drug Discovery Program", 21, False, TEAL)],
     [("on the Papulacandin / Fusacandin class (β-1,3-glucan / FKS1 inhibitors)", 16, False, RGBColor(0xC8, 0xD2, 0xDC))]])
txt(s, Inches(0.9), Inches(5.7), Inches(11.6), Inches(1.0),
    [[("Presenter:  ______________________          Date:  ____________", 14, False, RGBColor(0xC8, 0xD2, 0xDC))],
     [("A complete, reproducible pipeline — from scattered literature to a ready-to-test design.", 13, False, RGBColor(0x9C, 0xB4, 0xD8), True)]])
_page = 1

# =========================================================== 2 PROBLEM
s = header("Background · the problem", "Invasive fungal infections: an urgent, under-served threat")
chips = [("Rising incidence", "hundreds of thousands of deaths / year", CORAL),
         ("Very few drug classes", "limited options, serious side effects", AMBER),
         ("Resistance spreading", "existing drugs losing ground", BLUE)]
cw = Inches(3.87); x = ML; y = Inches(1.85)
for head_, sub, col in chips:
    card(s, x, y, cw, Inches(1.65))
    rect(s, x, y, Inches(0.14), Inches(1.65), col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, x + Inches(0.35), y + Inches(0.26), cw - Inches(0.5), Inches(1.2),
        [[(head_, 17, True, INK)], [(sub, 13, False, GREY)]])
    x += cw + Inches(0.26)
card(s, ML, Inches(3.95), CW, Inches(2.35), CARD2, line=None)
txt(s, ML + Inches(0.4), Inches(4.2), CW - Inches(0.8), Inches(1.95),
    [[("The opportunity — and the catch", 17, True, BLUE)],
     [("The ", 15, False, INK), ("papulacandin", 15, True, INK),
      (" natural products block the enzyme that builds the fungal cell wall (β-1,3-glucan "
       "synthase / FKS1). Humans have no cell wall → an attractive, low-toxicity target.", 15, False, INK)],
     [("❗  The catch: potent in a test tube, but ", 15, False, CORAL),
      ("activity is lost in blood serum", 15, True, CORAL),
      (" — so they never became drugs.", 15, False, CORAL)],
     [("If we understand and fix that, we revive an entire drug class.", 14, True, INK)]])

# =========================================================== EVIDENCE BASE (fig)
s = header("Background · the data", "The evidence base: decades of scattered literature, curated")
fig(s, "fig_dataset.png", ML, Inches(1.75), CW, Inches(3.95))
strip(s, [[("A clean, provenance-tracked benchmark — itself a lasting resource for the field.",
            15, True, WHITE)]], top=Inches(6.05))

# =========================================================== 3 SERUM GAP (fig)
s = header("The scientific question", "The serum gap, quantified: blood switches most analogs off")
fig(s, "fig_serum_gap.png", ML, Inches(1.6), CW, Inches(4.35))
strip(s, [[("“serum shift”  =  dose needed in serum  ÷  dose needed in broth", 16, True, WHITE),
           ("      — the one number we set out to lower.", 14, False, RGBColor(0xC8, 0xD2, 0xDC))]])

# =========================================================== 4 STRATEGY (fig)
s = header("Strategy", "A 14-stage funnel: start broad and cheap, narrow to confident designs")
fig(s, "fig_pipeline.png", ML, Inches(1.6), CW, Inches(4.5))
strip(s, [[("Each stage checks our reasoning before the next — and the whole platform is reusable.",
            15, True, WHITE)]])

# =========================================================== 5 METHODS (fig+text)
s = header("Methods · the algorithm", "The staged quantum-chemistry funnel")
fig(s, "fig_qm_funnel.png", ML, Inches(1.55), Inches(6.7), Inches(5.4))
tx = ML + Inches(7.0)
txt(s, tx, Inches(2.0), Inches(4.9), Inches(4.6),
    [[("How it works", 17, True, BLUE)],
     [("①  AI invents novel molecules.", 15, True, INK)],
     [("②  Cheap filters remove the obvious losers.", 15, False, INK)],
     [("③  Accurate quantum chemistry computes the real 3-D shapes — and re-checks the "
       "cheap results.", 15, False, INK)],
     [("④  The costliest method (DFT) touches only the few survivors.", 15, False, INK)],
     [("", 8, False, INK)],
     [("Same measurement engine at every tier → a fair, honest comparison.", 14, True, TEAL)]])

# =========================================================== 6 HONESTY LADDER (fig) money
s = header("Key methodological result", "The “honesty ladder”: more rigor shrank the easy signal")
fig(s, "fig_honesty_ladder.png", ML, Inches(1.55), CW, Inches(4.35))
strip(s, [[("A less careful analysis would have made a confident — and WRONG — claim at every rung. ",
            14, True, WHITE), ("Knowing what is NOT true is a real, durable result.", 14, False, RGBColor(0xC8, 0xD2, 0xDC))]])

# =========================================================== CONFOUND (fig)
s = header("Result · the decisive check", "Why the exciting signal collapsed: it was mostly potency")
fig(s, "fig_confound.png", ML, Inches(1.6), CW, Inches(4.45))
strip(s, [[("Serum outcome is governed by intrinsic potency (ρ=0.79); remove it and the 3-D signal vanishes — ",
            14, True, WHITE), ("so we model the potency-free SHIFT.", 14, False, RGBColor(0xC8, 0xD2, 0xDC))]])

# =========================================================== ELECTRONICS (fig+text)
s = header("Result · independent check", "A different property (electronics) — the same honest story")
fig(s, "fig_electronics.png", ML, Inches(1.6), Inches(7.5), Inches(5.0))
tx = ML + Inches(7.75)
txt(s, tx, Inches(2.05), Inches(4.3), Inches(4.4),
    [[("What it shows", 16, True, BLUE)],
     [("•  Polarizability looks like the strongest signal — but it is a ", 13.5, False, INK),
      ("size / potency proxy", 13.5, True, INK), (" (vanishes on the shift).", 13.5, False, INK)],
     [("•  ", 13.5, False, INK), ("QM logP", 13.5, True, TEAL),
      (" independently corroborates the polar-surface lead (+0.30 on the shift).", 13.5, False, INK)],
     [("", 8, False, INK)],
     [("Two unrelated descriptor families converge on the same weak, honest conclusion.",
       13.5, True, INK)]])

# =========================================================== 7 ECHINOCANDIN (fig+text)
s = header("Learning from approved drugs", "Echinocandins: same target, real clinical data")
fig(s, "fig_echinocandin.png", ML, Inches(1.55), Inches(6.9), Inches(5.5))
tx = ML + Inches(7.15)
txt(s, tx, Inches(1.95), Inches(4.75), Inches(1.0),
    [[("Two lessons that redirected our design:", 16, True, BLUE)]])
c1 = card(s, tx, Inches(2.75), Inches(4.75), Inches(1.5), LIGHT, line=LINE)
rect(s, tx, Inches(2.75), Inches(0.12), Inches(1.5), CORAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, tx + Inches(0.3), Inches(2.9), Inches(4.3), Inches(1.25),
    [[("1.  The tail is REQUIRED for potency", 14.5, True, INK)],
     [("can’t just make it water-friendly or cut it off, or the drug stops working.", 13, False, GREY)]])
c2 = card(s, tx, Inches(4.45), Inches(4.75), Inches(1.5), LIGHT, line=LINE)
rect(s, tx, Inches(4.45), Inches(0.12), Inches(1.5), GREEN, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, tx + Inches(0.3), Inches(4.6), Inches(4.3), Inches(1.25),
    [[("2.  SHAPE matters more than greasiness", 14.5, True, INK)],
     [("flexible tail = less serum loss. Our native tail is rigid & flat → the ‘bad’ zone.", 13, False, GREY)]])

# =========================================================== 8 AI DESIGN + REWARD (fig+text)
s = header("AI-guided design", "We built an AI generator — and validated its scoring on real data")
fig(s, os.path.join(OUT, "phase12_reward_validation.png"), ML, Inches(1.65), Inches(6.4), Inches(4.9))
tx = ML + Inches(6.8)
txt(s, tx, Inches(2.0), Inches(5.0), Inches(4.6),
    [[("The AI, done honestly", 17, True, BLUE)],
     [("•  An AI generator invents novel, on-scaffold molecules.", 14.5, False, INK)],
     [("•  Its ‘good design’ reward was ", 14.5, False, INK),
      ("checked against real measured data", 14.5, True, INK),
      (" (left) before we trusted it.", 14.5, False, INK)],
     [("•  It is built to produce a ", 14.5, False, INK),
      ("discriminating series", 14.5, True, TEAL),
      (" — molecules that TEST the idea, not just assume it.", 14.5, False, INK)],
     [("", 8, False, INK)],
     [("Left: the reward tracks true serum tolerance on the 24 known compounds "
       "(a weak-but-real trend, honestly reported).", 12.5, False, GREY, True)]])

# =========================================================== TAIL RANKING (fig)
s = header("Result · round 1 (detail)", "AI-designed tails, ranked by the cheap 3-D screen")
fig(s, "fig_tail_ranking.png", ML, Inches(1.55), CW, Inches(4.5))
strip(s, [[("Only 2 of 12 beat the native tail here — and the accurate re-check (next slide) removed even those.",
            14.5, True, WHITE)]])

# =========================================================== 9 ROUND-1 RESULT (fig)
s = header("Result · round 1", "The quantum funnel caught a false positive before synthesis")
fig(s, "fig_round1.png", ML, Inches(1.6), CW, Inches(4.55))
strip(s, [[("The accurate method exposed the cheap method’s ‘winner’ as an artifact — ", 14, True, WHITE),
           ("caught before any lab time was spent.", 14, False, RGBColor(0xC8, 0xD2, 0xDC))]])

# =========================================================== 10 SHORTLIST (fig+text)
s = header("The deliverable", "A ranked, easy-to-make shortlist — redirected by the drug lesson")
fig(s, "fig_tail_ladder.png", ML, Inches(1.55), Inches(7.0), Inches(4.6))
tx = ML + Inches(7.2)
txt(s, tx, Inches(1.85), Inches(4.7), Inches(0.6), [[("Make first (best bet):", 16, True, BLUE)]])
for i, (nm, sub, col) in enumerate([("Palmitoyl (saturated)", "cleanest single test", GREEN),
                                    ("Branched saturated", "caspofungin-like", GREEN),
                                    ("One-kink control", "make second", TEAL)]):
    yy = Inches(2.5) + i * Inches(0.82)
    card(s, tx, yy, Inches(4.7), Inches(0.72), RGBColor(0xE7, 0xF4, 0xEA) if col == GREEN else LIGHT, line=col)
    txt(s, tx + Inches(0.28), yy + Inches(0.05), Inches(4.3), Inches(0.6),
        [[(nm, 14, True, INK), (f"   ·  {sub}", 12.5, False, GREY)]], anchor=MSO_ANCHOR.MIDDLE)
strip(s, [[("Happy coincidence:  ", 15, True, TEAL),
           ("these are also the EASIEST to synthesize", 15, True, WHITE),
           (" — cheap fatty acids, one step. Best science = cheapest chemistry.", 14, False, RGBColor(0xC8, 0xD2, 0xDC))]])

# =========================================================== 11 CONTRIBUTIONS (grid)
s = header("Contributions", "What this program delivered — a reusable platform, not a one-off")
cards = [("📊  Benchmark dataset", "138 compounds, 1,042 records,\nfull provenance — a field resource"),
         ("🔬  Validated QM method", "multi-tier quantum pipeline +\nin-house 3-D surface engine"),
         ("🤖  AI molecule generator", "novel-structure design with a\nreward checked against real data"),
         ("💊  Cross-drug read-across", "clinical echinocandin evidence\nimported into the design logic"),
         ("⚙️  Reproducible engineering", "14 documented stages, automated\ntests, cluster run protocols"),
         ("🗺️  Roadmap + shortlist", "the exact experiment that\nunlocks predictive AI")]
cols = 3; gap = Inches(0.25); cw = (CW - gap * (cols - 1)) / cols; ch = Inches(2.05)
for i, (head_, body) in enumerate(cards):
    r, c = divmod(i, cols)
    x = ML + c * (cw + gap); y = Inches(1.9) + r * (ch + Inches(0.22))
    card(s, x, y, cw, ch, LIGHT, line=LINE)
    rect(s, x, y, cw, Pt(4), BLUE)
    txt(s, x + Inches(0.25), y + Inches(0.22), cw - Inches(0.5), ch - Inches(0.4),
        [[(head_, 15, True, INK)], [(body, 12.5, False, GREY)]])

# =========================================================== 12 LESSON + ROADMAP
s = header("The key lesson & the road ahead", "Where AI helps — and the loop that gets us there")
card(s, ML, Inches(1.85), CW, Inches(1.5), CARD2, line=None)
txt(s, ML + Inches(0.4), Inches(2.0), CW - Inches(0.8), Inches(1.25),
    [[("This is NOT “AI failed.”  ", 15, True, BLUE), ("It is a precise statement of the critical path:", 15, False, INK)],
     [("AI drug design needs a measured ‘answer key’. Physics-based guesses alone were shown — rigorously — to mislead.", 14, False, INK)],
     [("We now know exactly which few experiments unlock the AI — and the AI is already built to use them.", 14, True, GREEN)]])
loop = [("1.  Synthesize the\n4-molecule shortlist", BLUE), ("2.  Measure serum shift\n+ protein binding", AMBER),
        ("3.  Train the AI\non the results", GREEN), ("4.  AI designs\nround 2 (wider)", TEAL)]
n = len(loop); gap = Inches(0.55); cw = (CW - gap * (n - 1)) / n; y = Inches(4.1); x = ML
for i, (t, col) in enumerate(loop):
    card(s, x, y, cw, Inches(1.5), LIGHT, line=col)
    rect(s, x, y, Inches(0.12), Inches(1.5), col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, x + Inches(0.28), y, cw - Inches(0.4), Inches(1.5), [[(t, 13.5, True, INK)]], anchor=MSO_ANCHOR.MIDDLE)
    x += cw
    if i < n - 1:
        chevron(s, x + Inches(0.12), y + Inches(0.55), w=Inches(0.3), h=Inches(0.4), c=col)
    x += gap
txt(s, ML, Inches(5.85), CW, Inches(0.5),
    [[("↻  A virtuous loop: each round of data makes the AI a better designer.  ", 13.5, True, TEAL),
      ("The immediate next step needs no more computing.", 13.5, False, GREY)]], align=PP_ALIGN.CENTER)

# =========================================================== 13 SIGNIFICANCE
s = header("Significance", "Why this matters — and the one-line takeaway")
pts = [("Revived a stalled, high-value drug class as a tractable design problem", "urgent unmet clinical need"),
       ("Built an end-to-end AI + quantum-chemistry discovery platform", "reusable across programs"),
       ("Demonstrated rigorous self-correction (the ‘honesty ladder’)", "trustworthy, independent science"),
       ("Integrated chemistry + machine learning + quantum physics + pharmacology", "genuinely interdisciplinary"),
       ("Defined the exact experiment that unlocks predictive AI", "de-risks future investment")]
y = Inches(1.85)
for head_, sub in pts:
    rect(s, ML, y + Inches(0.06), Inches(0.16), Inches(0.16), TEAL, shape=MSO_SHAPE.OVAL)
    txt(s, ML + Inches(0.35), y, CW - Inches(0.4), Inches(0.6),
        [[(head_ + "  ", 14.5, True, INK), ("— " + sub, 13, False, GREY, True)]])
    y += Inches(0.62)
card(s, ML, Inches(5.2), CW, Inches(1.5), INK, line=None)
txt(s, ML + Inches(0.45), Inches(5.35), CW - Inches(0.9), Inches(1.25),
    [[("“We set out to teach a computer to design a serum-tolerant antifungal. We built the platform, "
       "curated the data, imported lessons from approved drugs, and — most importantly — learned to tell "
       "a real signal from a flattering illusion.", 14, False, WHITE, True)],
     [("The computer is ready; it awaits one clean experiment, designed down to the exact four molecules.”",
       14, True, TEAL, True)]])

# =========================================================== 14 THANK YOU
s = prs.slides.add_slide(BLANK); _bg(s, INK)
rect(s, 0, 0, SW, Inches(0.22), BLUE); rect(s, 0, SH - Inches(0.22), SW, Inches(0.22), TEAL)
txt(s, Inches(0.9), Inches(2.7), Inches(11.6), Inches(1.5),
    [[("Thank you", 44, True, WHITE)], [("Questions & discussion welcome", 20, False, TEAL)]])
txt(s, Inches(0.9), Inches(6.2), Inches(11.6), Inches(0.6),
    [[("Full technical detail, code, and reproducible reports: project repository (Phases 1–14).",
       13, False, RGBColor(0x9C, 0xB4, 0xD8), True)]])

# =========================================================== SPEAKER NOTES
NOTES = [
    # 1 TITLE
    "[~30s] Good morning, and thank you for being here. Over the next fifteen minutes I'll "
    "walk you through a drug-discovery program that combines artificial intelligence with "
    "quantum chemistry to attack a stubborn problem in antifungal medicine. I'll keep the "
    "chemistry light and focus on the ideas and the reasoning. The one-sentence version is on "
    "the slide: we set out to design antifungal molecules that keep working in blood — and "
    "along the way we built a complete, reusable computational platform and learned exactly "
    "where AI can and cannot help. Let me start with why this matters.",
    # 2 PROBLEM
    "[~60s] Invasive fungal infections are a growing, under-appreciated threat — they kill "
    "hundreds of thousands of people a year, often patients whose immune systems are already "
    "compromised. Yet we have only a handful of antifungal drug classes, many with serious "
    "side effects, and resistance is spreading. So there is a real need for genuinely new "
    "structures. Our starting point is a natural-product family, the papulacandins. They hit a "
    "target that is attractive because it is in the fungus but not in us — the machinery that "
    "builds the fungal cell wall. The catch, in red: these molecules are potent in a test tube "
    "but lose activity in blood serum. That single flaw is why they never became drugs. If we "
    "can understand and fix it, we potentially revive an entire class.",
    # 3 EVIDENCE BASE
    "[~45s] Before any modeling we had to build the dataset. The results we needed were "
    "scattered across decades of old papers, so we hand-curated them into one clean, "
    "computer-readable database — 138 compounds and over a thousand activity measurements, each "
    "traceable to its source. The heart of it is the 24 compounds with the matched measurement "
    "we care about: activity with and without serum. Thirteen keep some activity; eleven are "
    "switched off completely. We later imported an independent set of data on approved drugs "
    "for cross-checking. I want to flag that this curated dataset is itself a lasting "
    "contribution — it is the benchmark everything else is measured against, and others can reuse it.",
    # 4 SERUM GAP
    "[~50s] This is the problem in one picture. Every horizontal line is one compound. The blue "
    "dot on the left is the dose needed to stop the fungus in clean broth — low, meaning potent. "
    "Follow the line right to the dose needed in serum. Green compounds keep reasonable activity; "
    "the red arrows are compounds pushed past our detection ceiling — serum has essentially "
    "abolished them. The gap between the two ends is what we named the 'serum shift': serum dose "
    "divided by broth dose. A shift near one is ideal; a large shift is the clinical failure. "
    "Everything downstream is about lowering that shift without giving up potency. So the whole "
    "project reduces to one number we are trying to design down.",
    # 5 STRATEGY
    "[~45s] Here is the overall strategy. Rather than one big experiment, we built a staged "
    "funnel across four acts and fourteen documented steps. We start broad and cheap — curating "
    "data and looking for simple patterns. Then we go three-dimensional with quantum chemistry. "
    "Then — and this turned out to be pivotal — we borrow lessons from approved drugs. And "
    "finally we let AI design molecules and stress-test them. The key discipline, shown by the "
    "narrowing shape, is that each stage checks our reasoning before we spend effort on the next. "
    "And because it is all documented and reproducible, the platform outlives this one question.",
    # 6 QM FUNNEL
    "[~55s] This is the engine at the heart of Act Two — the 'algorithm', in plain terms. These "
    "molecules are floppy; think of a piece of cooked spaghetti that curls into many shapes. To "
    "judge them fairly we compute the real 3-D shapes they adopt in water and measure how much "
    "greasy versus water-friendly surface they expose. The trick is cost: each tier down the "
    "funnel is more accurate and more expensive. So we generate many molecules, filter cheaply, "
    "and spend the costly quantum methods only on the few survivors — using the same measurement "
    "engine at every tier so comparisons stay fair. This funnel is reusable infrastructure, not "
    "specific to this one question.",
    # 7 HONESTY LADDER
    "[~70s — the key slide] This is the most important methodological slide and the strongest "
    "evidence of how we work. The vertical axis is the strength of the apparent signal linking "
    "our 3-D descriptor to serum tolerance; read left to right as we increase rigor. Crude 2-D "
    "chemistry: weak. A cheap 3-D method: it looks strong — genuinely tempting, you could write a "
    "paper on that bar. But accurate quantum chemistry shrank it, and once we did the honest "
    "thing and controlled for plain potency, it collapsed to essentially zero. The exciting early "
    "result was largely an illusion. The message: at every rung, a less careful analysis would "
    "have made a confident and wrong claim. We climbed the whole ladder and reported what was "
    "actually there. Knowing precisely what is not true is a real, durable result — and it is "
    "what stops us wasting years chasing a mirage. "
    "[If asked 'so you found nothing?': we found what is NOT true — which saves years — plus one "
    "honest lead and a validated platform. Rigorously established negatives are results.]",
    # 8 CONFOUND
    "[~55s] Here is the data behind that collapse — the decisive check. On the left, the dose "
    "needed in serum plotted against the dose needed in broth. They track together tightly, "
    "correlation 0.79. In plain terms, a compound's fate in serum is mostly explained by how "
    "potent it was to begin with, not by a special serum-resistance property. On the right, once "
    "we statistically remove that potency effect, the 3-D descriptor that looked so promising — "
    "the greasy-surface one — drops to essentially zero. Only one weak signal survives: exposing "
    "water-friendly rather than greasy surface. That survivor is faint, but it is the honest lead "
    "we carried forward. This is what separates a real effect from a potency artifact.",
    # 9 ELECTRONICS
    "[~45s] We didn't stop at shape. We also examined a completely different property — the "
    "molecules' electronic structure — as an independent check, and got the same story. The "
    "descriptor that looks strongest here, polarizability, turns out to be another stand-in for "
    "size and potency: it vanishes on the potency-free shift. But one quantity — a "
    "quantum-computed measure of water-versus-oil preference — independently points the same way "
    "as our surface lead. Two unrelated families of measurement converging on the same modest "
    "conclusion. When independent methods agree you trust the direction, even while being honest "
    "that the effect is small.",
    # 10 ECHINOCANDINS
    "[~65s] Here is where borrowing from approved drugs pays off. The echinocandins are "
    "antifungals already on the market that hit the same target and carry the same kind of greasy "
    "tail, so their well-documented serum behavior is a natural second opinion. Two lessons. "
    "First: the tail is required for the drug to work at all — it anchors the molecule into the "
    "membrane — so we cannot simply make it water-friendly or cut it off. Second, and more "
    "useful: among drugs that keep their potency, the one with the smallest serum problem, "
    "caspofungin, has a flexible, branched tail, while the worst offenders have rigid, flat, "
    "aromatic tails — good zone lower-left, bad zone upper-right. Our native molecule's tail is a "
    "rigid, flat chain, sitting squarely in the bad zone. That single observation redirected our "
    "whole design.",
    # 11 AI GENERATOR
    "[~55s] Now the AI. We built a generator that invents new molecules on our scaffold. But — "
    "discipline again — we did not let it optimize a made-up target. We first checked its scoring "
    "rule against real measured data; that is the plot on the left, showing the reward genuinely "
    "tracks serum tolerance across our known compounds, honestly reported as a weak-but-real "
    "trend. And we designed it to produce a spread of molecules that would test our idea, not "
    "just confirm it. So the AI here is not a black box making bold claims — it is a disciplined "
    "design tool whose scoring we validated before trusting. "
    "[If asked 'why not a bigger neural network?': with no measured labels yet, a bigger model "
    "just learns our unvalidated guess faster — garbage in. That is why the next step is the lab.]",
    # 12 TAIL RANKING
    "[~45s] Round one: we asked the AI to redesign just the greasy tail, keeping the rest of the "
    "molecule fixed — a clean, one-variable experiment. This shows all twelve designed tails "
    "ranked by how much greasy surface they expose, lower being better, against the native "
    "baseline in red. At this cheaper tier only two beat the native tail — the two green bars. "
    "Encouraging. But we had learned not to trust the cheap tier, so we did the accurate "
    "re-check — the next slide.",
    # 13 ROUND-1 ARTIFACT
    "[~55s] And here is the payoff of the funnel discipline. On the left, twelve designs down to "
    "two, and then, at the most accurate quantum level, zero survivors. On the right is why. Our "
    "top candidate looked better than native with the cheaper method — but the accurate method "
    "showed its water-friendly group folds back and hides inside the molecule, leaving it "
    "actually worse than where we started. A computational artifact. The important part: we "
    "caught it before a single day of laboratory synthesis. That is exactly what the expensive "
    "quantum tier is for — cheap insurance against making the wrong molecule.",
    # 14 SHORTLIST
    "[~55s] So computation told us the water-friendly route was a dead end — and the "
    "approved-drug lesson told us the right route instead: don't make the tail water-friendly, "
    "make it less rigid. This is our deliverable: a short, ranked set of molecules that keep the "
    "tail the same length — so potency is preserved — and change only its shape, from rigid to "
    "flexible. And here is the practical gift: these are also the easiest molecules to make — "
    "cheap, off-the-shelf building blocks, one step — unlike the original or the exotic designs. "
    "So the best scientific bet and the cheapest chemistry are the same molecules. We make the "
    "flexible ones first.",
    # 15 CONTRIBUTIONS
    "[~50s] Stepping back, here is what the program delivered — and I want to stress it is a "
    "platform, not a one-off. A curated benchmark dataset that is a resource for the field. A "
    "validated quantum-chemistry method for these large, floppy natural products, which is "
    "genuinely hard. An AI generator with a scoring rule checked against real data. A framework "
    "for importing clinical drug knowledge into design. Reproducible engineering — documented "
    "stages, automated tests, run scripts — so anyone can rebuild it. And a clear roadmap with a "
    "concrete shortlist. Any one of these outlives the specific question we started with.",
    # 16 WHERE AI HELPS + ROADMAP
    "[~55s] I want to be direct about the honest conclusion, because it is often misunderstood. "
    "This is not 'AI failed.' It is a precise statement of the critical path: AI drug design "
    "needs a measured answer key to learn from, and we showed — rigorously — that physics-based "
    "guesses alone can mislead. The good news is we now know exactly which few experiments unlock "
    "the predictive AI, and the AI is already built to use them. The loop at the bottom is the "
    "path: synthesize the four-molecule shortlist, measure their serum behavior, feed that back "
    "to train the AI, and let it design the next, wider round. Each turn makes the AI a better "
    "designer — and the immediate next step needs no more computing, just a lab experiment.",
    # 17 SIGNIFICANCE
    "[~55s] To summarize the significance. We took a stalled but high-value drug class and turned "
    "it into a tractable design problem. We built an end-to-end AI-plus-quantum-chemistry "
    "platform that is reusable. We demonstrated rigorous self-correction — the honesty ladder — "
    "which is really the hallmark of trustworthy science. The work spans chemistry, machine "
    "learning, quantum physics, and clinical pharmacology. And we defined the exact experiment "
    "that unlocks predictive AI, which de-risks future investment. If I leave you with one line: "
    "we set out to teach a computer to design a serum-tolerant antifungal — and along the way we "
    "learned to tell a real signal from a flattering illusion. The computer is ready; it is "
    "waiting for one clean experiment, narrowed to exactly four molecules.",
    # 18 THANK YOU
    "[~20s] Thank you. I am happy to take questions — about the biology, the computational "
    "methods, the AI, or where this goes next. All the detail, the code, and the reproducible "
    "reports live in the project repository, so anything I show can be independently checked.",
]
for _sl, _note in zip(prs.slides, NOTES):
    _sl.notes_slide.notes_text_frame.text = _note

path = os.path.join(OUT, "Papulacandin_promotion_report.pptx")
prs.save(path)
print(f"Wrote {path}  ({len(prs.slides._sldIdLst)} slides, {len(NOTES)} speaker notes)")
