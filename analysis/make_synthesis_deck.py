#!/usr/bin/env python3
"""
make_synthesis_deck.py
======================
Render outputs/SYNTHESIS_phases1-10.md into a PowerPoint deck
(outputs/Papulacandin_serum_gap_synthesis.pptx). Self-contained: title +
bullet layouts built by hand, key phase figures embedded where they exist.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(HERE, "outputs")

NAVY = RGBColor(0x1F, 0x3A, 0x5F)
TEAL = RGBColor(0x2A, 0x9D, 0x8F)
CORAL = RGBColor(0xE7, 0x6F, 0x51)
GREY = RGBColor(0x55, 0x5B, 0x66)
LIGHT = RGBColor(0xF4, 0xF6, 0xF8)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _bg(slide, color=RGBColor(0xFF, 0xFF, 0xFF)):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def _bar(slide, color=NAVY, h=Inches(0.12)):
    s = slide.shapes.add_shape(1, 0, 0, SW, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()


def title_slide(title, subtitle, tag):
    s = prs.slides.add_slide(BLANK); _bg(s, NAVY)
    tf = _box(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(2.2))
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p2 = tf.add_paragraph(); p2.text = subtitle
    p2.font.size = Pt(20); p2.font.color.rgb = TEAL
    tf3 = _box(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.7))
    p3 = tf3.paragraphs[0]; p3.text = tag
    p3.font.size = Pt(13); p3.font.color.rgb = RGBColor(0xC8, 0xD2, 0xDC); p3.font.italic = True
    return s


def content_slide(title, bullets, accent=NAVY):
    """bullets: list of (level, text, bold, color-or-None)."""
    s = prs.slides.add_slide(BLANK); _bg(s)
    _bar(s, accent)
    th = _box(s, Inches(0.6), Inches(0.32), Inches(12.1), Inches(0.9))
    tp = th.paragraphs[0]; tp.text = title
    tp.font.size = Pt(28); tp.font.bold = True; tp.font.color.rgb = accent
    tf = _box(s, Inches(0.75), Inches(1.5), Inches(11.9), Inches(5.6))
    first = True
    for level, text, bold, color in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = ("• " if level == 0 else "   – ") + text
        p.level = level
        p.font.size = Pt(20 - 2 * level)
        p.font.bold = bool(bold)
        p.font.color.rgb = color or (NAVY if level == 0 else GREY)
        p.space_after = Pt(7)
    return s


def figure_slide(title, img, caption, accent=NAVY, bullets=None):
    s = prs.slides.add_slide(BLANK); _bg(s)
    _bar(s, accent)
    th = _box(s, Inches(0.6), Inches(0.32), Inches(12.1), Inches(0.9))
    tp = th.paragraphs[0]; tp.text = title
    tp.font.size = Pt(26); tp.font.bold = True; tp.font.color.rgb = accent
    path = os.path.join(OUT, img)
    img_w = Inches(7.4); max_h = Inches(5.4)
    if os.path.exists(path):
        pic = s.shapes.add_picture(path, Inches(0.5), Inches(1.5), width=img_w)
        if pic.height > max_h:                      # too tall: rescale both dims by one factor
            scale = max_h / pic.height
            pic.width = int(pic.width * scale)
            pic.height = int(max_h)
    # right-hand commentary
    tf = _box(s, Inches(8.2), Inches(1.6), Inches(4.7), Inches(5.2))
    pts = bullets or [(0, caption, False, None)]
    first = True
    for level, text, bold, color in pts:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = ("• " if level == 0 else "   – ") + text
        p.level = level
        p.font.size = Pt(16 - level)
        p.font.bold = bool(bold)
        p.font.color.rgb = color or (NAVY if level == 0 else GREY)
        p.space_after = Pt(8)
    return s


# ----- 1. title -----
title_slide(
    "Why do some papulacandins survive serum?",
    "A 10-phase computational dissection of the FKS1-inhibitor serum gap",
    "Papulacandin / fusacandin serum-gap study  ·  Phases 1–10  ·  generated from SYNTHESIS_phases1-10.md")

# ----- 2. bottom line up front -----
content_slide("Bottom line up front", [
    (0, "Serum POTENCY is governed by intrinsic potency × molecular size / lipophilicity", True, None),
    (0, "Serum TOLERANCE (the free→serum MIC shift) has, at most, a weak association with", True, None),
    (1, "exposing POLAR rather than HYDROPHOBIC surface", True, TEAL),
    (1, "|ρ| ≈ 0.30–0.33, p ≈ 0.13 — consistent across independent methods, never significant", False, None),
    (0, "No computable descriptor or docking model predicts serum tolerance on this dataset", True, CORAL),
    (0, "The binding constraint is the DATA (n=24, 11/24 censored), not the methods", True, None),
], accent=TEAL)

# ----- 3. the dataset -----
content_slide("The dataset — the dependent variable", [
    (0, "24 compounds with matched serum-free AND serum whole-cell MIC vs C. albicans", False, None),
    (1, "Yeung-1996 fusacandin analogs 6a–6u + Fusacandin A/B + Papulacandin B", False, None),
    (0, "13/24 retain measurable serum activity; 11/24 serum MICs CENSORED at 100", False, None),
    (1, "values tied on {12.5, 25, 50, 100} → this is the ceiling on every correlation", False, CORAL),
    (0, "Two endpoints used throughout:", True, None),
    (1, "raw serum MIC — but it is ~potency-dominated (ρ=0.79 with serum-free MIC)", False, None),
    (1, "serum SHIFT fold (serum/free) — the potency-independent measure of true tolerance", False, TEAL),
])

# ----- 4. the funnel -----
content_slide("Approach — a multi-tier descriptor funnel", [
    (0, "Phases 1–4 : curate serum-gap table + interpretable 2D design-rule score", False, None),
    (0, "Phase 5     : scaffold-constrained generative design → 12 novel candidates", False, None),
    (0, "Phase 6     : QM funnel — CREST ensembles, in-house SASA, Gaussian DFT I/O", False, None),
    (0, "Phases 7–8 : 3D shape retrospective (MMFF proxy → real CREST) + confound test", False, None),
    (0, "Phase 9     : GFN2-xTB electronic / solvation descriptors", False, None),
    (0, "Phase 10   : explicit human serum albumin docking", False, None),
    (0, "Each tier asks: does THIS descriptor family predict the serum shift?", True, NAVY),
])

# ----- 5. phases 1-5 -----
content_slide("Phases 1–5 — 2D fails, and the chemotype is bRo5", [
    (0, "Best 2D predictor (rigid-aromatic + polar design score) vs serum MIC:", False, None),
    (1, "Spearman ρ = 0.32, not significant", True, CORAL),
    (0, "Generative design cleaves/re-esterifies the C-6′ aromatic ester → 12 novel analogs", False, None),
    (0, "Whole chemotype is beyond-Rule-of-5:", True, None),
    (1, "QED ≈ 0.01–0.03, 0/30 pass Lipinski → 2D drug-likeness is uninformative", False, None),
    (1, "→ motivated the 3D / QM approach", False, TEAL),
])

# ----- 6. phase 7 -----
figure_slide("Phase 7 — a promising 3D lead (MMFF proxy)",
             "phase7_retrospective_qm.png", "",
             accent=TEAL,
             bullets=[
                 (0, "Fast RDKit-MMFF ensembles on the 24 knowns", False, None),
                 (0, "Hydrophobic SASA vs serum MIC:", True, None),
                 (1, "ρ = −0.45 (p = 0.029) — beats 2D", True, TEAL),
                 (0, "Promising enough to justify real CREST", False, None),
                 (0, "Caveat: proxy ensembles, to be confirmed", False, GREY),
             ])

# ----- 7. phase 8 -----
figure_slide("Phase 8 — real CREST overturns it: a potency artifact",
             "phase8_retrospective_crest.png", "",
             accent=CORAL,
             bullets=[
                 (0, "Real CREST/GFN-FF: signal shrinks to ρ=−0.31 (p=0.14)", False, None),
                 (0, "Decisive confound check:", True, None),
                 (1, "serum-free MIC alone → serum MIC ρ=0.79", False, None),
                 (1, "control potency → hydrophobic SASA ρ=0.02", True, CORAL),
                 (0, "Rigidity hypothesis dies (shift ρ≈−0.1)", False, None),
                 (0, "Only coherent lead: polar SASA vs shift ρ=−0.33", True, TEAL),
             ])

# ----- 8. phase 9 -----
figure_slide("Phase 9 — electronics converge on the polar-surface lead",
             "phase9_electronic.png", "",
             accent=TEAL,
             bullets=[
                 (0, "GFN2-xTB: dipole, gap, polarizability, QM logP", False, None),
                 (0, "Polarizability ρ=−0.54 (p=0.01) — but a SIZE proxy", False, None),
                 (1, "shift ρ=−0.02; tracks Rg, potency", False, GREY),
                 (0, "QM logP ↔ hydrophobic SASA (ρ=+0.76);", False, None),
                 (1, "shift ρ=+0.30 — corroborates polar-surface lead", True, TEAL),
                 (0, "Two descriptor families converge; neither significant", False, None),
             ])

# ----- 9. phase 10 -----
figure_slide("Phase 10 — explicit HSA docking is null",
             "phase10_hsa.png", "",
             accent=CORAL,
             bullets=[
                 (0, "Rigid ensemble surface docking (flexible = intractable at ~38 torsions)", False, None),
                 (0, "HSA binding vs serum shift:", True, None),
                 (1, "ρ = +0.22 (p=0.30) — sign OPPOSITE to sequestration", True, CORAL),
                 (0, "Flat under potency- and size-controls", False, None),
                 (0, "Not a size proxy (ρ=−0.15 vs polarizability)", False, GREY),
                 (0, "Caveat: Vina unfit for 1200 Da amphiphiles → cannot EXCLUDE sequestration", False, GREY),
             ])

# ----- 10. convergent finding -----
content_slide("The one convergent, defensible finding", [
    (0, "Three independent families agree on a single qualitative rule:", True, NAVY),
    (1, "3D shape (polar SASA), QM solvation (logP), weakly the dipole", False, None),
    (0, "“Among equipotent analogs, bias the exposed surface toward POLAR / H-bonding", True, TEAL),
    (0, "  character to reduce serum loss.”  |ρ| ≈ 0.30–0.33, p ≈ 0.13 — directional, not proven", True, TEAL),
    (0, "Rigorously established NEGATIVES:", True, CORAL),
    (1, "raw serum MIC is potency-dominated (ρ=0.79) — model the shift instead", False, None),
    (1, "the rigidity hypothesis is not supported", False, None),
    (1, "size / lipophilicity (clogP, Rg, polarizability) tracks potency, not tolerance", False, None),
    (1, "explicit HSA drug-site docking carries no signal", False, None),
], accent=TEAL)

# ----- 11. design output -----
content_slide("Design output (carried forward, with caveats)", [
    (0, "12 novel candidates taken through the full funnel; 3 finalists by multi-objective score:", False, None),
    (1, "cand01 quinolinecarbonyl · cand02 naphthoyl-6-OH · cand03 pyridylphenyl", True, NAVY),
    (1, "GFN2-reranked ensembles + ready Gaussian DFT inputs", False, None),
    (0, "Ranking unchanged, but the RATIONALE shifts after Phases 8–10:", True, None),
    (1, "prefer added polar / H-bonding surface on the variable C-6′ group,", False, TEAL),
    (1, "not just rigid aromatic bulk — a hypothesis to test, not a validated predictor", False, GREY),
])

# ----- 12. limitations + next steps -----
content_slide("Limitations → next steps are EXPERIMENTAL", [
    (0, "Limitations", True, CORAL),
    (1, "n=24, 11 censored, single chemotype/lab — caps achievable significance", False, None),
    (1, "MMFF proxy OVERSTATED effects (−0.45→−0.31) — trust QM-ensemble numbers", False, None),
    (1, "Vina unsuited to 1000–1200 Da amphiphiles — Phase-10 null can't exclude sequestration", False, None),
    (0, "Recommended next steps", True, TEAL),
    (1, "1. Acquire UNCENSORED serum MICs on more analogs / a fresh chemotype (biggest lever)", False, None),
    (1, "2. Direct HSA-binding assay (equilibrium dialysis / fluorescence displacement)", False, None),
    (1, "3. If in silico: model fatty-acid sites & other carriers; run finalist DFT (exploratory)", False, None),
])

# ----- 13. closing -----
title_slide(
    "A clean, fully-documented stopping point",
    "Computation mapped the serum gap and its limits; the lever is now experimental data",
    "Per-phase write-ups: outputs/phase7–10_findings.md  ·  full report: outputs/SYNTHESIS_phases1-10.md")

out = os.path.join(OUT, "Papulacandin_serum_gap_synthesis.pptx")
prs.save(out)
print(f"Wrote {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
