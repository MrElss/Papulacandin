#!/usr/bin/env python3
"""
make_full_deck.py
=================
Comprehensive Papulacandin serum-gap deck (Phases 1-10): every phase, the key
figures and native tables, and DETAILED speaker notes on every slide. Companion
asset ZIP (export_assets.py) carries standalone figures/tables in case a viewer
mis-renders embedded objects.
"""
import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(HERE, "outputs")

NAVY = RGBColor(0x1F, 0x3A, 0x5F); TEAL = RGBColor(0x2A, 0x9D, 0x8F)
CORAL = RGBColor(0xE7, 0x6F, 0x51); GREY = RGBColor(0x55, 0x5B, 0x66)
GOLD = RGBColor(0xE9, 0xC4, 0x6A); WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text.strip()


def _bg(slide, color=WHITE):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = color


def _bar(slide, color=NAVY):
    s = slide.shapes.add_shape(1, 0, 0, SW, Inches(0.12))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()


def _title(slide, text, accent=NAVY, size=26):
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.28), Inches(12.2), Inches(0.95))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(size); p.font.bold = True; p.font.color.rgb = accent


def _fmt(v):
    if isinstance(v, float):
        return f"{v:.3g}"
    s = str(v)
    return s if len(s) <= 40 else s[:37] + "…"


def bullets(slide, items, left=0.75, top=1.5, width=11.9, height=5.5, base=20):
    tf = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height)).text_frame
    tf.word_wrap = True
    first = True
    for it in items:
        lvl, txt, bold, color = (it + (False, None))[:4] if len(it) < 4 else it
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.text = ("• " if lvl == 0 else ("   – " if lvl == 1 else "       · ")) + txt
        p.level = lvl
        p.font.size = Pt(base - 2 * lvl)
        p.font.bold = bool(bold)
        p.font.color.rgb = color or (NAVY if lvl == 0 else GREY)
        p.space_after = Pt(6)


def add_table(slide, df, left, top, width, fontsize=11, header=NAVY, col_w=None):
    rows, cols = df.shape
    gt = slide.shapes.add_table(rows + 1, cols, Inches(left), Inches(top),
                                Inches(width), Inches(0.3 * (rows + 1))).table
    for j, c in enumerate(df.columns):
        cell = gt.cell(0, j); cell.text = str(c)
        pr = cell.text_frame.paragraphs[0]; pr.font.size = Pt(fontsize); pr.font.bold = True
        pr.font.color.rgb = WHITE; cell.fill.solid(); cell.fill.fore_color.rgb = header
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i in range(rows):
        for j in range(cols):
            cell = gt.cell(i + 1, j); cell.text = _fmt(df.iat[i, j])
            pr = cell.text_frame.paragraphs[0]; pr.font.size = Pt(fontsize - 1)
            pr.font.color.rgb = GREY
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xF1, 0xF4, 0xF8) if i % 2 else WHITE
    if col_w:
        for j, w in enumerate(col_w):
            gt.columns[j].width = Inches(w)
    return gt


def add_fig(slide, img, left=0.5, top=1.45, width=7.4, max_h=5.5):
    path = os.path.join(OUT, img)
    if not os.path.exists(path):
        return None
    pic = slide.shapes.add_picture(path, Inches(left), Inches(top), width=Inches(width))
    if pic.height > Inches(max_h):
        sc = Inches(max_h) / pic.height; pic.width = int(pic.width * sc); pic.height = int(Inches(max_h))
    return pic


def slide_blank(accent=NAVY):
    s = prs.slides.add_slide(BLANK); _bg(s); _bar(s, accent); return s


# ============================ TITLE ============================
s = prs.slides.add_slide(BLANK); _bg(s, NAVY)
tf = s.shapes.add_textbox(Inches(0.9), Inches(1.9), Inches(11.5), Inches(2.6)).text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "The Papulacandin Serum Gap"
p.font.size = Pt(46); p.font.bold = True; p.font.color.rgb = WHITE
p2 = tf.add_paragraph(); p2.text = "A ten-phase computational dissection of why some FKS1 inhibitors survive serum — and what to do next"
p2.font.size = Pt(21); p2.font.color.rgb = TEAL
p3 = s.shapes.add_textbox(Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.8)).text_frame
pp = p3.paragraphs[0]; pp.text = "Data curation · 2D SAR · generative design · CREST/QM ensembles · electronic descriptors · HSA docking"
pp.font.size = Pt(13); pp.font.color.rgb = RGBColor(0xC8, 0xD2, 0xDC); pp.font.italic = True
_notes(s, """
Welcome. This deck walks through a complete computational investigation of the 'serum gap' in
papulacandin/fusacandin-class FKS1 (1,3-beta-glucan synthase) inhibitors: potent antifungals in
buffer that lose activity in serum. Over ten phases I built the dataset, exhausted 2D structure-activity
analysis, designed novel analogs, ran QM conformer ensembles (CREST), computed 3D and electronic
descriptors, and finally docked the compounds to human serum albumin. The honest headline — which I'll
support with figures and statistics throughout — is that no computable descriptor cleanly predicts
serum tolerance on this dataset, but the methods converge on one defensible design direction and, more
importantly, they precisely define what experiment is needed next. Every figure and table here is also
provided as a separate ZIP so you can reuse them independently of this file.
""")

# ============================ EXEC SUMMARY ============================
s = slide_blank(TEAL); _title(s, "Executive summary — what this project established", TEAL)
bullets(s, [
    (0, "Serum POTENCY is dominated by intrinsic potency × molecular size / lipophilicity", True, None),
    (0, "Serum TOLERANCE (the free→serum MIC shift) is, at most, weakly associated with", True, None),
    (1, "exposing POLAR rather than HYDROPHOBIC surface  (|ρ|≈0.30–0.33, p≈0.13)", True, TEAL),
    (1, "the same direction appears in 3D shape AND independent QM-solvation descriptors", False, None),
    (0, "Rigorously falsified along the way:", True, CORAL),
    (1, "the 'rigidity → tolerance' hypothesis; the hydrophobic-SASA 'hit' (a potency artifact)", False, None),
    (1, "explicit HSA drug-site docking carries no signal", False, None),
    (0, "The binding constraint is the DATA (n=24, 11/24 serum MICs censored), not the methods", True, NAVY),
    (0, "→ the next decisive step is experimental, and is now precisely specified", True, GOLD),
], base=18)
_notes(s, """
This is the whole story in one slide. Two things are well-supported. First, raw serum MIC is largely a
restatement of intrinsic potency and size — bigger, more lipophilic, more potent molecules look better in
serum simply because they started better. Second, when you remove that potency confound and look at the
true serum *shift*, the only signal that survives — weakly — is that exposing polar rather than hydrophobic
surface helps. Crucially that same direction shows up independently in 3D shape descriptors (polar SASA) and
in quantum-chemical solvation descriptors (a QM logP), which is more convincing than any single number.
Equally valuable are the clean negatives: I falsified the original rigidity hypothesis, showed the early
hydrophobic-SASA result was an artifact of intrinsic potency, and found explicit albumin docking gives nothing.
None of the descriptors reach statistical significance, and I'll show that this is a hard ceiling imposed by a
small, heavily-censored dataset — which is exactly why the recommendation at the end is an experiment, not
another calculation.
""")

# ============================ OUTLINE ============================
s = slide_blank(NAVY); _title(s, "Contents")
bullets(s, [
    (0, "Background — invasive fungal disease, the FKS1 target, the serum gap", False, None),
    (0, "The dataset — 24 matched serum-free / serum MIC pairs", False, None),
    (0, "Phases 1–4 — 2D structure-activity & machine-learning baselines", False, None),
    (0, "Phases 5–6 — generative design & the QM (CREST) descriptor funnel", False, None),
    (0, "Phases 7–8 — 3D retrospective: a promising lead, then the confound that kills it", False, None),
    (0, "Phase 9 — quantum-chemical electronic / solvation descriptors", False, None),
    (0, "Phase 10 — explicit human serum albumin docking", False, None),
    (0, "Synthesis — the convergent finding, the negatives, and guidance for next research", False, GOLD),
])
_notes(s, """
The deck follows the actual research arc. We start with why these molecules matter and define the serum gap
quantitatively. Phases 1–4 are the classical 2D and machine-learning attempts and establish the baseline to
beat. Phases 5–6 build novel candidates and the quantum-mechanical infrastructure. Phases 7–10 are the heart:
a 3D retrospective that looked promising then failed a confound test, electronic descriptors that corroborated
a weak lead, and a direct docking test of the albumin-sequestration mechanism. We close with the synthesis and
detailed, prioritized guidance for follow-on work. I'll flag the single most important slide — the Phase-8
confound analysis — when we reach it.
""")

# ============================ SECTION: BACKGROUND ============================
s = slide_blank(NAVY); _title(s, "Background — why papulacandins, and what is the serum gap?", NAVY)
bullets(s, [
    (0, "Invasive fungal infections kill ~1.5 M people/year; resistance to existing classes is rising", False, None),
    (0, "FKS1 / (1,3)-β-D-glucan synthase is the validated echinocandin target (fungal cell wall)", False, None),
    (0, "Papulacandins & fusacandins are non-peptide glycolipid FKS1 inhibitors — a distinct chemotype", False, None),
    (1, "potent against the enzyme and in serum-free whole-cell assays", False, None),
    (0, "THE PROBLEM (\"serum gap\"): most lose activity in the presence of serum", True, CORAL),
    (1, "e.g. Papulacandin B: serum-free ~1.7 µM → serum ~111 µM (≈64× loss)", False, None),
    (1, "a few analogs retain serum activity — those positives are the entire SAR signal", False, TEAL),
    (0, "Goal: find a computable rule to DESIGN serum-tolerant analogs", True, NAVY),
])
_notes(s, """
Context for why this is worth a ten-phase effort. Invasive fungal disease is a large and growing cause of death,
and the echinocandin drugs that hit the FKS1 glucan-synthase target are one of only a few usable classes, so
new chemotypes against the same target are valuable. Papulacandins and fusacandins are a structurally distinct,
non-peptide glycolipid class that inhibit FKS1 well — but they suffer the 'serum gap': strong activity in buffer
that largely disappears in serum, presumably because serum proteins bind and sequester the drug. The number to
anchor on is Papulacandin B: about 1.7 micromolar without serum, about 111 with serum, a ~64-fold loss. A minority
of analogs buck this trend and stay active in serum; because they're so few, they carry essentially all of the
structure-activity information. The project's goal was to turn those few positives into a predictive, computable
design rule.
""")

# ============================ DATASET ============================
s = slide_blank(NAVY); _title(s, "The dataset — the dependent variable (24 matched pairs)")
dfp = pd.read_csv(os.path.join(OUT, "serum_gap_pairs.csv"))
sub = dfp[["compound_id", "name", "serumfree_mic_ugml", "serum_mic_ugml", "serum_shift_fold", "serum_active"]].copy()
sub.columns = ["ID", "name", "free MIC", "serum MIC", "shift×", "serum-active"]
add_table(s, sub.head(11), left=0.6, top=1.5, width=8.5, fontsize=10,
          col_w=[1.3, 2.6, 1.1, 1.1, 1.0, 1.4])
bullets(s, [
    (0, "24 compounds,", True, None), (0, "matched", False, None),
], left=9.3, top=1.5, width=3.6, base=14)
tf = s.shapes.add_textbox(Inches(9.25), Inches(1.5), Inches(3.7), Inches(5.4)).text_frame; tf.word_wrap = True
for i, (t, b, c) in enumerate([
        ("Yeung-1996 fusacandin series 6a–6u + Fusacandin A/B + Papulacandin B", False, GREY),
        ("13/24 keep measurable serum activity", True, TEAL),
        ("11/24 serum MICs CENSORED at 100", True, CORAL),
        ("values tied on {12.5,25,50,100}", False, GREY),
        ("Endpoints: serum MIC AND the serum SHIFT (serum/free) — the potency-free measure", False, NAVY)]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = "• " + t; p.font.size = Pt(13); p.font.bold = b; p.font.color.rgb = c; p.space_after = Pt(8)
_notes(s, """
This is the foundation everything rests on, and its limitations explain every result that follows. We have 24
compounds for which both a serum-free and a serum whole-cell MIC against Candida albicans exist — mostly the
Yeung 1996 fusacandin C-6-prime ester series, plus the reference natural products. The table shows the first 11;
the full file is in the asset pack. Two features dominate the statistics. First, 11 of 24 serum MICs are
right-censored at 100 — the assay simply reports '>100' or '=100' — and the measured values are tied on only four
levels, so any rank correlation is built on heavy ties and has a low ceiling. Second, throughout the project I use
two endpoints: the raw serum MIC, and the serum *shift fold* (serum divided by serum-free), which removes intrinsic
potency and is the cleaner measure of true serum tolerance. Keep the censoring in mind — it is the reason nothing
reaches significance later, no matter how good the descriptor.
""")

# ============================ STRATEGY ============================
s = slide_blank(NAVY); _title(s, "Strategy — a multi-tier descriptor funnel")
bullets(s, [
    (0, "Each phase asks one question: does THIS descriptor family predict the serum shift?", True, NAVY),
    (0, "Tier 1 — 2D topological (clogP, fsp3, MW, rotatable bonds, design-rule score)", False, None),
    (0, "Tier 2 — machine learning (external FKS model; within-series serum model)", False, None),
    (0, "Tier 3 — generative design → novel candidates for prospective testing", False, None),
    (0, "Tier 4 — 3D conformational: CREST ensembles → Boltzmann-weighted SASA / shape", False, None),
    (0, "Tier 5 — quantum-chemical electronics: GFN2-xTB dipole, gap, polarizability, QM logP", False, None),
    (0, "Tier 6 — explicit mechanism: docking to human serum albumin", False, None),
    (0, "Cheap, broad tiers screen; expensive QM tiers confirm — and guard against false leads", True, TEAL),
])
_notes(s, """
The design philosophy is a funnel: start with cheap, broad descriptors and only escalate to expensive quantum
chemistry where it's justified. Tier 1 is classical 2D medicinal-chemistry descriptors. Tier 2 is machine learning,
both an external FKS-inhibitor model and a within-series serum model. Tier 3 is generative design to produce novel
analogs. Tier 4 moves to true three-dimensional, conformer-ensemble descriptors computed from CREST quantum
conformational searches — solvent-accessible surface area split into hydrophobic and polar, plus shape. Tier 5 adds
electronic and solvation descriptors from semi-empirical quantum chemistry. Tier 6 tests the mechanism directly by
docking to albumin. A key methodological lesson, which we'll see concretely in Phases 7 and 8, is that the cheap
tiers can produce false leads, and the expensive QM tier is what catches them — so the funnel is as much about
discipline as efficiency.
""")

# ---- helper for figure+notes phase slides ----
def fig_phase(title, img, side, notes, accent=NAVY, fig_w=7.5):
    s = slide_blank(accent); _title(s, title, accent)
    add_fig(s, img, width=fig_w)
    tf = s.shapes.add_textbox(Inches(8.3), Inches(1.55), Inches(4.6), Inches(5.4)).text_frame; tf.word_wrap = True
    first = True
    for lvl, txt, bold, color in side:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.text = ("• " if lvl == 0 else "   – ") + txt
        p.font.size = Pt(15 - lvl); p.font.bold = bool(bold)
        p.font.color.rgb = color or (NAVY if lvl == 0 else GREY); p.space_after = Pt(7)
    _notes(s, notes)
    return s

# ============================ PHASE 1 ============================
fig_phase("Phase 1 — which properties separate tolerant vs killed?",
          "phase1_descriptor_boxplots.png",
          [(0, "Serum-tolerant analogs tend to higher MW, lower fsp3", False, None),
           (1, "i.e. larger, more rigid-aromatic", False, None),
           (0, "MW: ρ=−0.59 (p=0.002); fsp3: ρ=+0.58 (p=0.003)", True, TEAL),
           (0, "clogP does NOT separate them (ρ=−0.33, p=0.12)", True, CORAL),
           (0, "Hypothesis (n=24, one series): rigid-aromatic character, not bulk lipophilicity", False, None)],
          notes="""
Phase 1 is the first SAR pass. Splitting the 24 compounds into serum-tolerant and serum-killed and comparing
2D descriptors, the cleanest separators are molecular weight (tolerant compounds are heavier) and fraction sp3
carbon (tolerant compounds are less sp3, i.e. more aromatic/rigid), each significant at p≈0.002–0.003 by both
Spearman and Mann-Whitney. Importantly, clogP — the obvious lipophilicity descriptor you'd expect to drive serum
protein binding — does NOT separate the groups. That is the first hint that the naive 'lipophilic compounds get
sequestered' story is too simple. This motivated the original working hypothesis: rigid aromatic character, rather
than bulk lipophilicity, tracks serum tolerance. Note the caveat that will recur — this is one chemical series and
n=24, so these are hypotheses, not laws. The boxplots here are in the asset pack as phase1_descriptor_boxplots.""",
          accent=NAVY)

# ============================ PHASE 1 table ============================
s = slide_blank(NAVY); _title(s, "Phase 1 — descriptor statistics (tolerant vs killed)")
d1 = pd.read_csv(os.path.join(OUT, "phase1_descriptor_stats.csv"))
d1d = d1[["descriptor", "median_tolerant", "median_killed", "spearman_rho_vs_serumMIC", "spearman_p", "mannwhitney_p"]]
d1d.columns = ["descriptor", "median (tol)", "median (killed)", "ρ vs serum MIC", "Spearman p", "MW-U p"]
add_table(s, d1d, left=1.2, top=1.7, width=10.8, fontsize=13)
bullets(s, [(0, "MW and fsp3 are the significant separators; clogP, HBD, rotatable bonds are not", True, NAVY)],
        left=1.2, top=5.7, width=11, base=15)
_notes(s, """
Here are the Phase-1 numbers in full. Reading down the Spearman and Mann-Whitney p-value columns, only molecular
weight and fraction-sp3 clear significance; hydrogen-bond acceptors and TPSA are borderline; and clogP, HBD count,
and rotatable-bond count are flat. The medians make the direction concrete: tolerant compounds are ~56 Da heavier
and notably less sp3. The practical message is that the discriminating signal, such as it is, lives in size and
aromatic rigidity — not in the lipophilicity or flexibility descriptors a medicinal chemist would reach for first.
This table is phase1_descriptor_stats in the asset pack.""")

# ============================ PHASE 2 ============================
fig_phase("Phase 2 — where in the cascade do compounds fail?",
          "phase2_attrition.png",
          [(0, "Per-compound ladder: enzyme → serum-free cell → serum cell", False, None),
           (0, "Compounds are potent on the enzyme AND serum-free cells", True, TEAL),
           (0, "Activity is lost ONLY at the serum step", True, CORAL),
           (1, "localises the problem to serum exposure, not target affinity or cell entry", False, None),
           (0, "Confirms the serum gap is a bioavailability/sequestration problem", False, NAVY)],
          notes="""
Phase 2 asks a diagnostic question: at which step does activity disappear? By assembling a per-compound ladder —
purified FKS1 enzyme inhibition, then serum-free whole-cell MIC, then serum whole-cell MIC — we see that the
compounds are potent on the enzyme and on serum-free cells, and only collapse at the serum step. That cleanly
localizes the serum gap: it is not a target-affinity problem and not a cell-entry problem, it is specifically a
serum-exposure problem, consistent with serum-protein sequestration reducing free drug. This is what justifies
spending the rest of the project on serum-binding-related descriptors rather than on potency or permeability. The
attrition and activity-ladder figures (phase2_attrition, phase2_ladder_anchors) are in the asset pack.""",
          accent=NAVY)

# ============================ PHASE 3A ============================
fig_phase("Phase 3A — can an external FKS model rank papulacandins?",
          "fig_phase3a_chemspace.png",
          [(0, "Trained an external FKS-inhibitor classifier; mapped papulacandins into its space", False, None),
           (0, "CV AUC 0.997 — but on a trivial split", True, CORAL),
           (0, "Papulacandins are OUT-OF-DOMAIN (model scores ≤0.39)", True, NAVY),
           (0, "Honest negative: use the model as a chemical-space reference, not a predictor", False, GREY)],
          notes="""
Phase 3A tests whether an existing, externally-trained FKS-inhibitor machine-learning model can rank our compounds.
The cross-validated AUC looks spectacular at 0.997, but that is on a trivial actives-versus-decoys split and is not
meaningful for our question. When we project the papulacandins into the model's descriptor space — the chemical-space
map shown here — they sit outside the training distribution and all score low (≤0.39). The honest conclusion is a
negative: the external model is out of domain for this chemotype and should be used only as a chemical-space reference,
not as a serum or activity predictor. I'm including the negative deliberately, because knowing a tool does NOT apply
saved us from over-trusting it. Figures: fig_phase3a_chemspace, fig_phase3a_cv_bars, fig_phase3a_score_dist.""",
          accent=CORAL)

# ============================ PHASE 3B ============================
fig_phase("Phase 3B — can we predict serum MIC within the series?",
          "phase3b_pred_vs_obs.png",
          [(0, "Within-series regularized model for serum MIC", False, None),
           (0, "In-sample looks good: ρ=0.65", False, None),
           (0, "Leave-one-out collapses: ρ=0.14 (n.s.)", True, CORAL),
           (1, "classic small-n over-fitting (optimism)", False, None),
           (0, "Verdict: not yet predictive — data-limited", True, NAVY)],
          notes="""
Phase 3B builds a model specifically for serum MIC within the fusacandin series. In-sample it correlates at Spearman
0.65, which is tempting — but the honest test is leave-one-out cross-validation, and there the correlation collapses
to 0.14, not significant. The gap between in-sample 0.65 and leave-one-out 0.14 is textbook small-sample optimism:
with 24 compounds and many descriptors, a model can memorize but not generalize. The optimism figure quantifies that
shrinkage. The verdict is that we cannot yet build a predictive serum model from this data — not because the chemistry
is unlearnable, but because the dataset is too small and too censored. This is the moment the project pivots from
fitting models to understanding mechanism with physics-based descriptors. Figures: fig_phase3b_coeffs,
fig_phase3b_optimism, phase3b_pred_vs_obs.""",
          accent=CORAL)

# ============================ PHASE 3 SUMMARY TABLE ============================
s = slide_blank(NAVY); _title(s, "Phase 3 — summary of the modeling attempts")
d3 = pd.read_csv(os.path.join(OUT, "table_phase3_summary.csv"))
add_table(s, d3, left=0.5, top=1.7, width=12.3, fontsize=11, col_w=[0.8, 4.0, 5.3, 2.2])
_notes(s, """
This table consolidates the modeling phases. Phase 1 gave a hypothesis — rigid-aromatic character tracks tolerance,
clogP doesn't — caveated by n=24 and one series. Phase 2 gave a robust per-compound ladder localizing failure to the
serum step. Phase 3A was an honest negative: the external FKS model is out of domain. Phase 3B was data-limited:
in-sample 0.65 but leave-one-out 0.14. The throughline is that classical 2D and ML approaches extract a weak hypothesis
but no predictive model, and they keep pointing at the same ceiling — the data. That is the explicit justification for
escalating to the quantum-mechanical 3D descriptors in the later phases. This is table_phase3_summary in the pack.""")

# ============================ PHASE 4 ============================
fig_phase("Phase 4 — existing serum-tolerant leads & design rules",
          "fig_phase4_existing_leads.png",
          [(0, "Catalogued the 13 serum-active analogs as SAR anchors", False, None),
           (0, "Best leads (serum MIC 12.5): 6j, 6k, 6n, 6o, 6p", True, TEAL),
           (1, "all extended biphenyl / terphenyl C-6′ aroyl esters", False, None),
           (0, "Distilled into a transparent, source-attributed design score", False, None),
           (0, "Used to rank/triage — NOT a black box", False, NAVY)],
          notes="""
Phase 4 turns the positives into design guidance. The 13 serum-active analogs are the only positive examples we have,
so they are the SAR anchors. The best of them — fusacandin analogs 6j, 6k, 6n, 6o, 6p, all at serum MIC 12.5 — share
an extended biphenyl or terphenyl aromatic ester at the C-6-prime position. From these I distilled an interpretable,
source-attributed 'design-rule score' rewarding rigid extended aromatics plus a polar handle, deliberately transparent
rather than a black box so a chemist can see why a compound scores well. This score becomes the 2D baseline that the
later 3D work has to beat. Figure fig_phase4_existing_leads and tables table_phase4_existing_leads /
table_phase4_candidate_ranking are in the pack; the rules themselves are in design_rules.md.""",
          accent=NAVY)

# ============================ PHASE 4 leads table ============================
s = slide_blank(NAVY); _title(s, "Phase 4 — the serum-tolerant leads (SAR anchors)")
d4 = pd.read_csv(os.path.join(OUT, "table_phase4_existing_leads.csv"))
d4d = d4[["name", "serumfree_mic_ugml", "serum_mic_ugml", "serum_shift_fold", "modification_summary"]].head(8).copy()
d4d.columns = ["analog", "free MIC", "serum MIC", "shift×", "C-6′ modification"]
add_table(s, d4d, left=0.5, top=1.6, width=12.3, fontsize=11, col_w=[2.0, 1.2, 1.2, 1.0, 6.9])
_notes(s, """
A closer look at the leads. Every top serum-tolerant analog carries an extended aromatic ester at C-6-prime —
biphenyl, terphenyl, alkoxy-biphenyl variants — on the conserved fusacandin core, and they all start from the same
strong serum-free MIC of 0.78. The 'shift' column is the fold loss in serum; the best leads still shift 16-fold but
retain usable serum activity because they began so potent. This table previews a theme that becomes decisive in
Phase 8: because the conserved core is identical across analogs, the serum differences must come from this variable
C-6-prime group, and the leads are simultaneously the most potent and the most serum-tolerant — a confound we will
have to disentangle. Table: table_phase4_existing_leads.""")

# ============================ PHASE 5 retrospective (2D baseline) ============================
fig_phase("Phases 4→5 — the 2D baseline to beat",
          "phase5_retrospective.png",
          [(0, "2D design-rule score vs observed serum MIC", False, None),
           (0, "Spearman ρ = 0.32 (p≈0.20) — not significant", True, CORAL),
           (0, "Best the 2D world can do", False, None),
           (0, "Sets the bar for every 3D / QM descriptor that follows", True, NAVY)],
          notes="""
This is the number to remember: the best 2D design-rule score correlates with observed serum MIC at Spearman 0.32,
not significant. This is the baseline every subsequent, more expensive descriptor has to beat to justify itself.
It also frames the central tension of the project — 2D is weakly suggestive but not predictive — and motivates the
generative and quantum-mechanical work. Whenever you see a rho on a later slide, compare it back to this 0.32.
Figure: phase5_retrospective.""",
          accent=CORAL)

# ============================ PHASE 5 generative ============================
fig_phase("Phase 5 — generative design of novel analogs",
          "phase5_score_distributions.png",
          [(0, "Cleave & re-esterify the validated C-6′ aromatic ester (InChIKey-verified)", False, None),
           (0, "37 acyl fragments (data-derived + designed) → multi-objective scoring", False, None),
           (1, "QED, synthetic accessibility, design rules, novelty — no learned activity model", False, None),
           (0, "12 novel candidates advanced; honest given the tiny dataset", True, TEAL)],
          notes="""
Phase 5 designs new molecules. Using the validated chemistry — cleaving the aromatic C-6-prime ester and
re-esterifying with new acyl groups, verified by InChIKey round-trip — I generated a virtual library from 37 acyl
fragments, both data-derived and rationally designed (rigid aromatics, some with polar handles). Each candidate is
scored by a transparent multi-objective function: drug-likeness (QED), synthetic accessibility, the Phase-1 design
rules, and novelty versus known compounds. Deliberately, there is NO learned activity predictor in the loop — with
24 noisy datapoints that would be dishonest. Twelve novel candidates advanced to the QM funnel. The score
distributions are shown here; the full 30-row library is phase5_virtual_library in the pack.""",
          accent=NAVY)

# ============================ PHASE 5 bRo5 ============================
s = slide_blank(CORAL); _title(s, "Phase 5 — a key cheminformatics finding: the chemotype is bRo5", CORAL)
bullets(s, [
    (0, "Across all 30 designs: QED ≈ 0.01–0.03 and 0/30 pass Lipinski's Rule of 5", True, CORAL),
    (0, "These are 'beyond-Rule-of-5' (bRo5) glycolipids — like cyclosporine, not like aspirin", False, None),
    (0, "Consequence: standard 2D drug-likeness metrics are SATURATED and uninformative here", True, NAVY),
    (1, "QED, Lipinski, etc. cannot rank within this chemotype — they're all ~equally 'bad'", False, None),
    (0, "→ This is the formal justification for moving to 3D / conformational descriptors", True, TEAL),
    (1, "if 2D can't even differentiate the molecules, only 3D shape/electronics can", False, None),
], base=19)
_notes(s, """
A short but important cheminformatics result. Every one of the 30 designed analogs has a near-zero QED and fails
Lipinski's Rule of Five — these are 'beyond-Rule-of-5' molecules, the same category as cyclosporine and other large
natural-product-like therapeutics. The practical consequence is that standard 2D drug-likeness metrics are saturated:
they rate every compound in this chemotype as equally 'undruglike,' so they carry no discriminating information within
the series. This isn't a flaw in the molecules; it's a flaw in applying small-molecule heuristics to a large-molecule
class. It is the formal justification for the entire second half of the project: if 2D descriptors cannot even tell
these molecules apart, then any real signal must live in their three-dimensional shape and electronics, which is what
the QM funnel computes.""")

# ============================ PHASE 6 ============================
s = slide_blank(NAVY); _title(s, "Phase 6 — the QM descriptor funnel (CREST → 3D descriptors)")
bullets(s, [
    (0, "CREST conformer ensembles (GFN-FF, water) for each compound — hundreds of conformers", False, None),
    (0, "In-house vectorized Shrake–Rupley SASA (no external dependency), Boltzmann-weighted", False, None),
    (1, "split into HYDROPHOBIC vs POLAR exposed surface — the serum-binding-relevant quantity", False, None),
    (0, "Ensemble SASA spread = a genuine 3D flexibility/rigidity metric", False, None),
    (0, "Radius of gyration, asphericity = overall shape", False, None),
    (0, "Gaussian DFT inputs (B3LYP/6-31G(d), PCM water, ESP charges) for populated conformers", False, None),
    (0, "Self-tested end-to-end before any real QM data existed", True, TEAL),
], base=18)
_notes(s, """
Phase 6 is the engineering that makes the 3D analysis possible. For each compound we run a CREST quantum
conformational search to get a realistic Boltzmann ensemble of hundreds of conformers in implicit water, then compute
descriptors on the whole ensemble rather than a single guessed structure. The central descriptor is solvent-accessible
surface area, computed with an in-house vectorized Shrake-Rupley algorithm so there's no fragile external dependency,
and crucially split into hydrophobic versus polar exposed surface — that split is the physically-motivated proxy for
serum-protein binding. We also compute the ensemble spread of hydrophobic SASA as a true 3D flexibility metric, plus
radius of gyration and asphericity for shape, and we generate Gaussian DFT inputs for the populated conformers to get
electrostatics later. The pipeline was self-tested end to end on synthetic ensembles before any real cluster data
arrived. The 12-candidate descriptor table is phase6_qm_descriptors in the pack.""")

# ============================ PHASE 7 ============================
fig_phase("Phase 7 — a promising 3D lead (MMFF-proxy ensembles)",
          "phase7_retrospective_qm.png",
          [(0, "Fast RDKit-MMFF ensembles on the 24 knowns (proxy for CREST)", False, None),
           (0, "Boltzmann-weighted HYDROPHOBIC SASA vs serum MIC:", False, None),
           (1, "ρ = −0.45 (p = 0.029) — beats the 2D baseline of 0.32", True, TEAL),
           (0, "First descriptor to reach significance — promising enough to commit cluster time", True, NAVY),
           (0, "Caveat flagged at the time: proxy ensembles, must confirm with real CREST", False, GREY)],
          notes="""
Phase 7 is the first apparent win. Before committing expensive cluster CREST time, I ran a fast proxy — RDKit MMFF
conformer ensembles — on the 24 known compounds and computed the same hydrophobic-SASA descriptor. It correlates with
serum MIC at Spearman −0.45, p=0.029: the first descriptor in the whole project to beat the 2D baseline of 0.32 and to
reach significance. The negative sign means more exposed hydrophobic surface goes with lower (better) serum MIC, which
at face value is a real, usable lead. This was promising enough to justify running the real quantum CREST ensembles to
confirm it. But I want to flag that I marked it as provisional at the time precisely because MMFF is a proxy — and the
next slide shows why that caution mattered. Figure: phase7_retrospective_qm.""",
          accent=TEAL)

# ============================ PHASE 8 ============================
fig_phase("Phase 8 — real CREST overturns the lead",
          "phase8_retrospective_crest.png",
          [(0, "Real CREST/GFN-FF ensembles (24 cmpds, 132–597 conformers each)", False, None),
           (0, "Hydrophobic SASA weakens to ρ=−0.31 (p=0.14) — loses significance", True, CORAL),
           (0, "Rigidity metric (SASA spread): flat (ρ≈+0.17)", False, None),
           (0, "MMFF proxy OVERSTATED the effect (−0.45 → −0.31)", True, NAVY),
           (0, "Methodological lesson: confirm proxy hits with real QM", True, TEAL)],
          notes="""
Phase 8 runs the real quantum CREST ensembles and re-tests Phase 7. The hydrophobic-SASA correlation weakens from
−0.45 to −0.31 and loses significance (p=0.14). The rigidity metric — the ensemble spread of hydrophobic SASA, which
was the original Phase-1 hypothesis operationalized in 3D — is essentially flat. So with better physics the headline
lead shrinks. The first lesson is methodological and general: a promising result from cheap proxy conformers must be
confirmed with real QM, because the proxy systematically overstated the effect. But the more important result is on
the next slide — it isn't just that the number got smaller; it's WHY, and that turns out to invalidate the lead
entirely. Figure: phase8_retrospective_crest.""",
          accent=CORAL)

# ============================ PHASE 8 confound (KEY) ============================
s = slide_blank(GOLD); _title(s, "Phase 8 — the decisive confound test  ★", RGBColor(0xB8, 0x86, 0x0B))
d8 = pd.read_csv(os.path.join(OUT, "phase8_confound_analysis.csv"))
d8d = d8.copy(); d8d.columns = ["descriptor", "partial ρ (serum MIC | potency)", "partial p", "ρ vs SHIFT", "shift p"]
add_table(s, d8d, left=0.6, top=1.55, width=12.1, fontsize=12)
bullets(s, [
    (0, "Serum-free MIC ALONE tracks serum MIC at ρ=0.79 — intrinsic potency dominates", True, CORAL),
    (0, "Control for potency → hydrophobic SASA collapses to ρ=0.02 (the Phase-7 'hit' was an artifact)", True, NAVY),
    (0, "Only survivor pointing the right way: POLAR SASA vs shift ρ=−0.33 (p=0.12)", True, TEAL),
], left=0.6, top=3.5, width=12.1, base=15)
_notes(s, """
This is the single most important slide in the deck — the star marks it. The question: is the hydrophobic-SASA
signal real, or is it just intrinsic potency in disguise? First fact: serum-free MIC by itself predicts serum MIC at
Spearman 0.79 — so most of what 'predicts serum MIC' is simply how potent the compound was to begin with, carried
through. Second, when you compute the PARTIAL correlation of hydrophobic SASA with serum MIC while controlling for
serum-free potency, it collapses from −0.31 to +0.02 — statistically nothing. In other words, the entire Phase-7 lead
was a potency artifact: bigger, more hydrophobic molecules happened to be more potent, and that potency, not any
serum-specific property, drove the correlation. The one descriptor that survives and points the mechanistically
sensible way is POLAR SASA against the pure serum shift, at −0.33 (p=0.12): more polar exposed surface, less serum
loss. It's not significant, but it is real-direction and not a potency artifact, so it becomes the lead we carry
forward. Table: phase8_confound_analysis.""")

# ============================ PHASE 8b ============================
s = slide_blank(NAVY); _title(s, "Phase 8b — GFN-FF → GFN2 refinement of the finalists")
d8b = pd.read_csv(os.path.join(OUT, "phase8b_gfnff_vs_gfn2.csv"))
cols = ["compound", "n_conformers__gfnff", "n_conformers__gfn2",
        "hydrophobic_sasa_mean__gfnff", "hydrophobic_sasa_mean__gfn2",
        "hydrophobic_sasa_std__gfnff", "hydrophobic_sasa_std__gfn2"]
d8bd = d8b[cols].copy()
d8bd.columns = ["finalist", "n(GFNFF)", "n(GFN2)", "hSASA FF", "hSASA GFN2", "hSASA sd FF", "hSASA sd GFN2"]
add_table(s, d8bd, left=0.7, top=1.7, width=11.9, fontsize=12)
bullets(s, [
    (0, "Re-ranking the 3 finalists' ensembles at GFN2 shifted populations materially", True, NAVY),
    (1, "cand02: a flat GFN-FF set → one dominant GFN2 conformer (68%) — DFT set would have been wrong", False, CORAL),
    (0, "Confirms: refine before the expensive DFT step", True, TEAL),
], left=0.7, top=4.0, width=11.9, base=15)
_notes(s, """
A methodological aside that doubles as quality control. For the three design finalists I re-ranked the cheap GFN-FF
conformer ensembles at the more accurate GFN2 level. The populations shifted materially — most strikingly cand02,
where a flat, many-conformer GFN-FF picture collapsed to a single dominant GFN2 conformer at 68% population. The
practical point: if we had generated the expensive DFT single points from the GFN-FF populations, we'd have computed
electronics on the wrong conformers. This validates the funnel discipline of refining conformer energies before the
most expensive step, and it's why the finalist DFT inputs were regenerated from GFN2. Table: phase8b_gfnff_vs_gfn2.""")

# ============================ PHASE 9 ============================
fig_phase("Phase 9 — quantum electronic / solvation descriptors",
          "phase9_electronic.png",
          [(0, "GFN2-xTB single points (water + octanol), 24 cmpds, 102 conformers", False, None),
           (0, "Polarizability α(0): biggest raw serum-MIC ρ=−0.54 (p=0.01)…", False, None),
           (1, "…but a SIZE proxy (shift ρ=−0.02; tracks Rg, potency)", True, CORAL),
           (0, "QM logP ↔ hydrophobic SASA (ρ=0.76); shift ρ=+0.30", True, TEAL),
           (1, "independently corroborates the polar-surface lead", False, None),
           (0, "Two descriptor families now converge — neither significant", True, NAVY)],
          notes="""
Phase 9 adds the electronics the shape descriptors miss, via cheap GFN2-xTB single points in water and octanol — about
200 calculations. Two results. First, molecular polarizability is the single largest raw correlate of serum MIC in the
whole project at −0.54, p=0.01 — but it is a size descriptor in disguise: it tracks radius of gyration and intrinsic
potency, and against the pure serum shift it is zero. So like hydrophobic SASA before it, it's a potency/size echo, not
a tolerance handle. Second, and more interesting, a quantum logP computed from the water/octanol solvation free-energy
difference correlates with hydrophobic SASA at 0.76 — the electronic and shape pictures are measuring the same physical
thing — and gives a serum-shift correlation of +0.30: more hydrophobic, worse tolerance. That mirrors the Phase-8 polar-
SASA lead of −0.33 in both direction and magnitude. So two independent descriptor families now converge on the same
qualitative rule, even though neither is individually significant. Figure: phase9_electronic; stats:
phase9_electronic_stats.""",
          accent=TEAL)

# ============================ PHASE 9 stats table ============================
s = slide_blank(NAVY); _title(s, "Phase 9 — electronic descriptor statistics")
d9 = pd.read_csv(os.path.join(OUT, "phase9_electronic_stats.csv"))
d9d = d9.copy(); d9d.columns = ["descriptor", "ρ serumMIC", "p", "partial ρ", "partial p", "ρ shift", "shift p"]
add_table(s, d9d, left=1.0, top=1.8, width=11.3, fontsize=13)
bullets(s, [(0, "logP_xtb is the descriptor that points the mechanistically-expected way on the SHIFT (+0.30)", True, TEAL),
            (0, "α(0) is big vs serum MIC but vanishes on the shift — a size/potency proxy", True, CORAL)],
        left=1.0, top=5.6, width=11, base=14)
_notes(s, """
The Phase-9 numbers in full. Scan the last two columns — correlation with the potency-independent serum shift — because
that's the endpoint that matters. logP_xtb is the standout there at +0.30 in the expected direction. Polarizability,
despite being the biggest number against raw serum MIC (−0.54), is essentially zero against the shift, confirming it as
a size/potency proxy. Dipole points weakly the right way; solvation free energy and HOMO-LUMO gap are flat. The
consistent message with Phase 8 is the convergence on the hydrophobic/polar-surface axis. Table: phase9_electronic_stats;
the per-compound descriptor values are phase9_electronic_descriptors.""")

# ============================ PHASE 10 ============================
fig_phase("Phase 10 — explicit human serum albumin docking",
          "phase10_hsa.png",
          [(0, "Direct test of the sequestration mechanism (HSA, PDB 1AO6)", False, None),
           (0, "Rigid ensemble surface docking (flexible = intractable at ~38 torsions)", False, None),
           (0, "HSA affinity vs serum shift: ρ=+0.22 (p=0.30) — NULL", True, CORAL),
           (1, "sign even opposite to sequestration; flat under potency/size controls", False, None),
           (0, "Not a size proxy (ρ=−0.15 vs α) → captured real surface binding, no signal", False, NAVY),
           (0, "Caveat: Vina unfit for 1200 Da amphiphiles → cannot EXCLUDE the mechanism", False, GREY)],
          notes="""
Phase 10 tests the mechanism head-on by docking every compound to human serum albumin. First I had to establish a valid
method: these are 1000–1200 Da glycolipids with about 38 rotatable bonds. Rigid docking into a drug pocket clashes badly
(the ligand is far bigger than the pocket), and fully flexible docking at 38 torsions is both intractable and unreliable.
So I used rigid ensemble surface docking — rigid albumin, our real QM conformers kept rigid, large boxes over both Sudlow
drug sites, best score over conformers and sites. The result is a clean null: HSA binding does not predict the serum shift
(rho=+0.22, p=0.30), and the sign is if anything opposite to what sequestration predicts. Reassuringly the score is not
just a size proxy, so the docking captured real surface association — it simply doesn't track tolerance. The honest caveat:
Vina's scoring function isn't parameterized for amphiphiles this large, and we modeled only the drug sites, not the fatty-
acid sites these lipid-tailed molecules might actually use — so a null here cannot exclude albumin sequestration; it points
to an experimental binding assay. Figure: phase10_hsa.""",
          accent=CORAL)

# ============================ SYNTHESIS TABLE ============================
s = slide_blank(GOLD); _title(s, "Synthesis — every approach, one table", RGBColor(0xB8, 0x86, 0x0B))
synth = pd.DataFrame([
    ["2D descriptors (Ph 1–5)", "design-rule score vs serum MIC", "ρ = 0.32", "n.s. — baseline"],
    ["ML serum model (Ph 3B)", "within-series, leave-one-out", "ρ = 0.14", "over-fit / data-limited"],
    ["3D hydrophobic SASA (Ph 7→8)", "MMFF proxy → real CREST", "−0.45 → −0.31", "potency artifact (partial 0.02)"],
    ["3D rigidity (SASA spread)", "ensemble flexibility", "≈ +0.17", "hypothesis falsified"],
    ["3D polar SASA (Ph 8)", "vs serum SHIFT", "−0.33 (p=0.12)", "lead — direction sensible"],
    ["QM polarizability (Ph 9)", "vs serum MIC / shift", "−0.54 / −0.02", "size/potency proxy"],
    ["QM logP (Ph 9)", "vs serum SHIFT", "+0.30 (p=0.15)", "corroborates polar lead"],
    ["HSA docking (Ph 10)", "vs serum SHIFT", "+0.22 (p=0.30)", "null"],
], columns=["approach", "test", "Spearman ρ", "verdict"])
add_table(s, synth, left=0.5, top=1.55, width=12.3, fontsize=12, col_w=[3.2, 3.4, 2.2, 3.5])
_notes(s, """
The entire project on one slide. Read top to bottom: 2D descriptors gave the 0.32 baseline; the within-series ML model
over-fit to a leave-one-out 0.14; the 3D hydrophobic-SASA lead looked good at −0.45 but shrank to −0.31 and then proved
to be a potency artifact under the confound test; the rigidity hypothesis was falsified; polar SASA against the serum
shift survived as a −0.33 directional lead; quantum polarizability was a big-but-empty size proxy; quantum logP
independently corroborated the polar lead at +0.30; and explicit albumin docking was null. The shape of the table is the
result: every method agrees that raw serum potency is a size/potency story, and the only thing pointing at genuine serum
tolerance — across two independent descriptor families — is exposing polar rather than hydrophobic surface, never quite
reaching significance. That convergence, plus the clean negatives, is the scientific output.""")

# ============================ CONVERGENT FINDING ============================
s = slide_blank(TEAL); _title(s, "The one convergent, defensible finding", TEAL)
bullets(s, [
    (0, "Among EQUIPOTENT analogs, biasing exposed surface toward POLAR / H-bonding character", True, TEAL),
    (0, "is associated with smaller serum loss.", True, TEAL),
    (1, "polar SASA vs shift ρ=−0.33; QM logP vs shift +0.30; dipole weakly aligned", False, None),
    (1, "two independent descriptor families (3D shape + QM solvation) agree", False, NAVY),
    (0, "Effect size |ρ|≈0.30–0.33, p≈0.12–0.15 — DIRECTIONAL, not proven", True, CORAL),
    (0, "Actionable design hypothesis for the variable C-6′ group:", True, NAVY),
    (1, "add polar / hydrogen-bonding surface, don't just extend rigid aromatic bulk", False, None),
], base=19)
_notes(s, """
If you take one actionable idea from the project, it is this. Among compounds of similar intrinsic potency, biasing the
molecule's exposed surface toward polar, hydrogen-bonding character — rather than hydrophobic surface — is associated with
less activity lost in serum. This is supported by polar SASA at −0.33 against the shift, a quantum logP at +0.30 in the
same direction, and a weakly-aligned dipole, i.e. two physically independent descriptor families agreeing. I'm careful to
call it a directional design hypothesis, not a validated predictor, because the effect size is modest and the p-values are
around 0.12–0.15. Concretely, for the variable C-6-prime aroyl group, it argues for adding polar or H-bonding substituents
rather than simply extending rigid aromatic bulk. It's the kind of hypothesis a focused synthetic series could confirm or
kill quickly.""")

# ============================ NEGATIVES ============================
s = slide_blank(CORAL); _title(s, "Equally valuable — what we rigorously ruled out", CORAL)
bullets(s, [
    (0, "Raw serum MIC is the WRONG endpoint — ~potency-dominated (ρ=0.79); always model the shift", True, NAVY),
    (0, "The Phase-1 'rigidity → tolerance' hypothesis is NOT supported (shift ρ≈−0.1)", True, CORAL),
    (0, "Bulk lipophilicity / size (clogP, Rg, polarizability) tracks POTENCY, not tolerance", True, CORAL),
    (1, "the Phase-7 hydrophobic-SASA 'hit' was a potency artifact (partial ρ=0.02)", False, None),
    (0, "Explicit HSA drug-site docking carries NO signal (and can't be trusted to exclude the mechanism)", True, CORAL),
    (0, "Proxy (MMFF) ensembles OVERSTATE effects — trust QM-ensemble numbers only", True, NAVY),
], base=18)
_notes(s, """
Negative results are results, and these are hard-won. One: raw serum MIC is the wrong thing to optimize because it's
~80% intrinsic potency; always model the serum shift or covary out potency. Two: the original rigidity hypothesis that
launched the project does not survive 3D testing. Three: every bulk size/lipophilicity descriptor — clogP, radius of
gyration, polarizability — tracks potency rather than tolerance, and the seductive Phase-7 hydrophobic-SASA result was
specifically shown to be a potency artifact. Four: explicit docking to albumin's drug sites carries no signal, though we
're honest that the method can't fully exclude the mechanism. Five, methodological: cheap proxy conformer ensembles
overstate effects, so only QM-ensemble numbers should be trusted. Collectively these stop a future team from wasting
effort re-chasing leads that look attractive but are already falsified.""")

# ============================ DESIGN OUTPUT ============================
s = slide_blank(NAVY); _title(s, "Design output — 3 finalists through the full funnel")
bullets(s, [
    (0, "12 novel candidates built (Phase 5); top 3 by multi-objective score taken to GFN2 + DFT-ready:", False, None),
    (1, "cand01 quinolinecarbonyl · cand02 naphthoyl-6-OH · cand03 pyridylphenyl", True, NAVY),
    (1, "all extended rigid aromatics; GFN2-reranked ensembles; Gaussian DFT inputs prepared", False, None),
    (0, "Post-Phase-8/9 the ranking is unchanged but the RATIONALE shifts:", True, TEAL),
    (1, "favour added polar / H-bonding surface on the C-6′ group, not aromatic bulk alone", False, None),
    (0, "All inputs are cluster-ready (CREST, GFN2 rerank, Gaussian) and committed to the repo", False, GREY),
    (0, "Caveat: a hypothesis to test prospectively, not a validated predictor", True, CORAL),
], base=18)
_notes(s, """
The tangible design deliverable. Twelve novel candidates were generated and the top three by the transparent
multi-objective score — quinolinecarbonyl, naphthoyl-6-OH, and pyridylphenyl analogs — were carried all the way through
the funnel: CREST ensembles, GFN2 re-ranking, and ready-to-run Gaussian DFT inputs. Importantly, the Phase-8/9 findings
don't change which three rank highest, but they change WHY: rather than just rewarding extended rigid aromatic bulk, the
guidance now is to bias the variable C-6-prime group toward polar, hydrogen-bonding surface. Every input file is
cluster-ready and committed, so synthesis and further computation can proceed immediately. As always, these are
prospective hypotheses to test, not predictions to trust blindly.""")

# ============================ GUIDANCE 1 ============================
s = slide_blank(GOLD); _title(s, "Guidance for subsequent research — methodology", RGBColor(0xB8, 0x86, 0x0B))
bullets(s, [
    (0, "Model the serum SHIFT (or covary out serum-free MIC) — never raw serum MIC", True, NAVY),
    (0, "Treat fast (MMFF/2D) results as SCREENS; confirm any lead with QM ensembles before believing it", True, NAVY),
    (0, "Refine conformer populations (GFN2) before any expensive DFT/docking on top", False, None),
    (0, "For bRo5 glycolipids, drop 2D drug-likeness metrics — they're saturated; use 3D/QM", False, None),
    (0, "Report partial correlations against potency by default — it is the dominant confound here", True, TEAL),
    (0, "Keep negative results first-class — they redirected the project twice and saved compute", False, None),
], base=18)
_notes(s, """
Methodological guidance for whoever picks this up. First and most important: never optimize raw serum MIC; model the
serum shift or explicitly control for serum-free potency, because potency is the dominant confound and will manufacture
false leads otherwise. Second: treat any cheap 2D or MMFF result as a screen only, and confirm with real QM ensembles
before acting — we saw a −0.45 'hit' evaporate. Third: refine conformer energies at GFN2 before spending DFT or docking
on top, or you'll compute expensive properties on the wrong conformers. Fourth: for beyond-Rule-of-5 glycolipids, stop
using 2D drug-likeness metrics — they're saturated and uninformative — and go straight to 3D/QM. Fifth: report
potency-partialled correlations by default. And sixth, keep negative results first-class; here they redirected the
project twice and saved substantial compute.
""")

# ============================ GUIDANCE 2 ============================
s = slide_blank(GOLD); _title(s, "Guidance — the next decisive steps are EXPERIMENTAL", RGBColor(0xB8, 0x86, 0x0B))
bullets(s, [
    (0, "1. ACQUIRE DATA — the single biggest lever", True, TEAL),
    (1, "more analogs with UNCENSORED serum MICs; ideally a second chemotype", False, None),
    (1, "powers the polar-surface lead to significance (or kills it) — current n=24/11-censored caps everything", False, None),
    (0, "2. MEASURE the mechanism directly", True, TEAL),
    (1, "HSA-binding assay (equilibrium dialysis / fluorescence-probe displacement at Sudlow I & II)", False, None),
    (1, "include fatty-acid sites — the lipid tails may bind there, not the drug sites docking modeled", False, None),
    (0, "3. Targeted synthesis to test the lead", True, TEAL),
    (1, "matched pairs differing only in polar vs hydrophobic C-6′ surface, equipotent serum-free", False, None),
    (0, "4. If in silico: model FA sites / other carriers (AAG, lipoproteins); run the prepared finalist DFT", False, GREY),
], base=17)
_notes(s, """
The most important guidance: the next decisive steps are experimental, not computational, and the project has narrowed
them precisely. Priority one is data — more analogs with uncensored serum MICs, ideally including a second chemotype.
This is the single biggest lever, because the current n=24 with 11 censored values is a hard statistical ceiling; with
even 40–60 uncensored points the polar-surface lead could be pushed to significance or cleanly rejected. Priority two is
to measure the mechanism directly with an albumin-binding assay — equilibrium dialysis or fluorescence-probe displacement
at the Sudlow sites — and critically to include the fatty-acid binding sites, because these lipid-tailed molecules may
bind there rather than at the drug sites our docking covered. Priority three is a small, targeted synthetic series of
matched pairs that differ only in polar versus hydrophobic C-6-prime surface while holding serum-free potency constant —
the cleanest possible test of the one lead. Only after those would more in-silico work — fatty-acid-site modeling, other
serum carriers, the prepared finalist DFT single points — be worthwhile, and even then as exploratory support.
""")

# ============================ REPRO / ASSETS ============================
s = slide_blank(NAVY); _title(s, "Reproducibility & deliverables")
bullets(s, [
    (0, "All code, inputs, and outputs committed to the repository (main)", False, None),
    (1, "phase5_generate.py · phase6_qm_layer.py · phase7–10 scripts · gen_*/make_* utilities", False, None),
    (0, "CREST/GFN2/DFT/xtb/docking inputs are cluster-ready and version-controlled", False, None),
    (0, "Per-phase write-ups: phase7–10_findings.md; master narrative: SYNTHESIS_phases1-10.md", False, None),
    (0, "This deck: Papulacandin_serum_gap_synthesis_FULL.pptx (with full speaker notes)", True, TEAL),
    (0, "Companion asset pack: Papulacandin_deck_assets.zip", True, TEAL),
    (1, "18 figures + 21 tables (CSV + rendered PNG) + 8 reports + INDEX", False, None),
], base=18)
_notes(s, """
Everything here is reproducible and committed to the repository's main branch: the generative pipeline, the QM
descriptor layer, the Phase 7–10 analysis scripts, and all the input-generation utilities, plus every cluster-ready
input for CREST, GFN2, DFT, xtb, and docking. The reasoning is documented at two levels — per-phase findings files and
a single master synthesis narrative. This presentation, with the full speaker notes you're reading, is saved alongside a
companion asset ZIP containing all 18 figures and all 21 tables — each table as both raw CSV and a rendered PNG image so
you can drop them into any document if a viewer mis-renders an embedded object — plus the eight written reports and an
index. In short, a new team can pick this up and either run the next calculation or hand the experimental spec to a wet
lab without reverse-engineering anything.
""")

# ============================ CLOSING ============================
s = prs.slides.add_slide(BLANK); _bg(s, NAVY)
tf = s.shapes.add_textbox(Inches(0.9), Inches(2.4), Inches(11.5), Inches(2.6)).text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "Computation mapped the serum gap — and its limits"
p.font.size = Pt(34); p.font.bold = True; p.font.color.rgb = WHITE
p2 = tf.add_paragraph(); p2.text = ("Potency is size/lipophilicity; tolerance weakly favours polar surface; "
                                    "the lever is now experimental data.")
p2.font.size = Pt(19); p2.font.color.rgb = TEAL
p3 = s.shapes.add_textbox(Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.7)).text_frame
pp = p3.paragraphs[0]; pp.text = "Thank you — full report: outputs/SYNTHESIS_phases1-10.md  ·  assets: Papulacandin_deck_assets.zip"
pp.font.size = Pt(13); pp.font.color.rgb = RGBColor(0xC8, 0xD2, 0xDC); pp.font.italic = True
_notes(s, """
To close: this was a rigorous, ten-phase computational effort that did three things well. It built and characterized the
serum-gap dataset; it exhaustively tested whether 2D, 3D, electronic, and explicit-docking descriptors can predict serum
tolerance — and honestly reported that none reach significance on this data; and it extracted the one convergent,
mechanistically-sensible design lead (favour polar over hydrophobic exposed surface) while falsifying several attractive
but wrong hypotheses. The decisive insight is that the bottleneck is data, not methods, and the project now specifies
exactly which experiments would break the impasse. Thank you — the full written narrative and the complete figure/table
asset pack are referenced here for anyone continuing the work.
""")

out = os.path.join(OUT, "Papulacandin_serum_gap_synthesis_FULL.pptx")
prs.save(out)
print(f"Wrote {out}  ({len(prs.slides._sldIdLst)} slides, with speaker notes)")
